#!/usr/bin/env python3
"""
Build the FitXpress Product Marketing Playbook deck for Nika as a .pptx,
styled on DESIGN.md tokens and structured after the FitXpress Sales Playbook (June 2026).

Source of content: workspace/playbooks/fitxpress-marketing-playbook-nika.md
No figure appears here that is not in that file's approved proof set (section 7.3).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import math, sys

# ----------------------------------------------------------------------------- tokens
BLUE      = RGBColor(0x14, 0x3D, 0xFF)
NAVY      = RGBColor(0x05, 0x0F, 0x40)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)
BLUE_50   = RGBColor(0xEC, 0xEF, 0xFF)
BLUE_100  = RGBColor(0xD8, 0xDE, 0xFF)
BLUE_200  = RGBColor(0xB1, 0xBD, 0xFF)
BLUE_300  = RGBColor(0x8A, 0x9C, 0xFF)
BLUE_400  = RGBColor(0x4F, 0x6D, 0xFF)
BLUE_600  = RGBColor(0x0F, 0x2E, 0xCD)
BLUE_700  = RGBColor(0x0B, 0x22, 0x99)
BLUE_800  = RGBColor(0x08, 0x18, 0x6B)
G50       = RGBColor(0xF9, 0xF9, 0xF9)
G100      = RGBColor(0xF2, 0xF2, 0xF2)
G200      = RGBColor(0xE5, 0xE5, 0xE5)
G300      = RGBColor(0xD1, 0xD1, 0xD1)
G400      = RGBColor(0xA8, 0xA8, 0xA8)
G500      = RGBColor(0x80, 0x80, 0x80)
G600      = RGBColor(0x66, 0x66, 0x66)
G700      = RGBColor(0x4C, 0x4C, 0x4C)
INK       = RGBColor(0x1A, 0x1A, 0x1A)

FONT = "Satoshi"

SW, SH = 13.3333, 7.5
ML     = 0.72                 # left margin
CW     = SW - 2 * ML          # content width
FOOTER = "FitXpress Product Marketing Playbook  ·  August 2026  ·  Confidential, internal use only"

warnings = []


# ----------------------------------------------------------------------------- helpers
def set_radial_gradient(shape, stops, center=(50000, 24000)):
    """Radial (path=circle) gradient fill, since DESIGN.md specifies a navy radial glow."""
    spPr = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    gs = "".join(
        f'<a:gs pos="{p}"><a:srgbClr val="{c}"/></a:gs>' for p, c in stops
    )
    cx, cy = center
    xml = (
        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'rotWithShape="1">'
        f"<a:gsLst>{gs}</a:gsLst>"
        f'<a:path path="circle"><a:fillToRect l="{cx}" t="{cy}" '
        f'r="{100000-cx}" b="{100000-cy}"/></a:path>'
        "</a:gradFill>"
    )
    frag = etree.fromstring(xml)
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(frag)
    else:
        spPr.append(frag)


def no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, fill=None, radius=None, line=None, line_w=1.0):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        # adjustment is a fraction of half the shorter side
        try:
            s.adjustments[0] = min(0.5, radius / (min(w, h) * 72.0))
        except Exception:
            pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        no_line(s)
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


MIN_SCALE = 0.80          # never shrink body copy below 80% of its declared size


def _measure(paras, w, size, line_spacing, space_after, scale):
    """Estimated rendered height in inches at a given scale factor."""
    total = 0.0
    for p in paras:
        pieces = p if isinstance(p, list) else [(p, {})]
        chars = sum(len(t) for t, _ in pieces)
        psize = max([o.get("size", size) for _, o in pieces] or [size]) * scale
        cpl = max(8, (w * 72.0) / (psize * 0.485))
        lines = max(1, math.ceil(chars / cpl)) if chars else 1
        total += lines * psize * line_spacing / 72.0
        if space_after:
            total += space_after * scale / 72.0
    return total


def txt(slide, x, y, w, h, runs, size=11, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, line_spacing=1.28, space_after=0, fit_check=True,
        italic=False, caps=False):
    """runs: str, or list of paragraphs. A paragraph is a str or a list of (text, dict) tuples.

    Text is measured before rendering and the point size is scaled down (to a floor of
    MIN_SCALE) so copy stays inside its box. Anything that still does not fit is reported,
    because that is a copy-length problem rather than a layout problem.
    """
    paras = runs if isinstance(runs, list) else [runs]

    scale = 1.0
    if fit_check:
        while scale > MIN_SCALE and _measure(paras, w, size, line_spacing,
                                             space_after, scale) > h + 0.02:
            scale -= 0.02
        need = _measure(paras, w, size, line_spacing, space_after, scale)
        if need > h + 0.02:
            warnings.append(
                f"  OVERFLOW  y={y:.2f} w={w:.2f} h={h:.2f} need~{need:.2f}in "
                f"@{size * scale:.1f}pt :: {str(paras[0])[:60]!r}"
            )

    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        if space_after:
            para.space_after = Pt(space_after * scale)
        pieces = p if isinstance(p, list) else [(p, {})]
        for t, o in pieces:
            r = para.add_run()
            r.text = t
            f = r.font
            f.name = FONT
            f.size = Pt(o.get("size", size) * scale)
            f.bold = o.get("bold", bold)
            f.italic = o.get("italic", italic)
            f.color.rgb = o.get("color", color)
            if caps or o.get("caps"):
                r.font._rPr.set("cap", "all")
    return tb


def eyebrow(slide, label, color=BLUE, y=0.46):
    tb = slide.shapes.add_textbox(Inches(ML), Inches(y), Inches(CW), Inches(0.24))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = label
    f = r.font
    f.name, f.size, f.bold = FONT, Pt(10), True
    f.color.rgb = color
    # letter-spacing 0.14em
    r.font._rPr.set("spc", "170")
    return tb


def slide_head(slide, label, title, sub=None, dark=False, title_size=29):
    tc = WHITE if dark else NAVY
    sc = BLUE_200 if dark else G600
    ec = BLUE_300 if dark else BLUE
    eyebrow(slide, label, color=ec)
    txt(slide, ML, 0.76, CW, 0.62, title, size=title_size, color=tc, bold=True,
        line_spacing=1.06)
    if sub:
        txt(slide, ML, 1.42, CW - 0.3, 0.50, sub, size=12.5, color=sc, line_spacing=1.26)
        return 2.00
    return 1.56


def footer_line(slide, page, total, dark=False):
    c = BLUE_200 if dark else G400
    txt(slide, ML, SH - 0.44, CW - 1.1, 0.22, FOOTER, size=7.5, color=c, fit_check=False)
    txt(slide, SW - ML - 1.0, SH - 0.44, 1.0, 0.22, f"{page} / {total}", size=7.5,
        color=c, align=PP_ALIGN.RIGHT, fit_check=False)


def navy_bg(slide, center=(50000, 22000), grid=True):
    bg = rect(slide, -0.02, -0.02, SW + 0.04, SH + 0.04)
    set_radial_gradient(
        bg,
        [(0, "3A57F5"), (32000, "0B2299"), (66000, "08186B"), (100000, "050F40")],
        center=center,
    )
    if grid:
        # faint measurement-grid texture
        for i in range(1, 20):
            x = i * (SW / 20.0)
            ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(0),
                                        Emu(4000), Inches(SH))
            ln.fill.solid()
            ln.fill.fore_color.rgb = RGBColor(0x2A, 0x45, 0xC8)
            no_line(ln)
            ln.shadow.inherit = False
            _set_alpha(ln, 12000)
        for i in range(1, 11):
            y = i * (SH / 11.0)
            ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(y),
                                        Inches(SW), Emu(4000))
            ln.fill.solid()
            ln.fill.fore_color.rgb = RGBColor(0x2A, 0x45, 0xC8)
            no_line(ln)
            ln.shadow.inherit = False
            _set_alpha(ln, 12000)


def _set_alpha(shape, alpha):
    """alpha in 0..100000 where 100000 = opaque."""
    clr = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
    if clr is not None:
        a = etree.SubElement(clr, qn("a:alpha"))
        a.set("val", str(alpha))


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ----------------------------------------------------------------------------- blocks
def cards(slide, items, top, cols=3, gap=0.2, height=None, fill=G50, radius=20,
          tsize=11.5, bsize=8.8, line=None, accent_bar=False, tcolor=NAVY,
          bcolor=G700, pad=0.19):
    """items: list of (title, body). body may be a list of bullet strings."""
    n = len(items)
    rows = math.ceil(n / cols)
    cwid = (CW - gap * (cols - 1)) / cols
    if height is None:
        height = (SH - 0.62 - top - gap * (rows - 1)) / rows
    for i, (t, b) in enumerate(items):
        r, c = divmod(i, cols)
        x = ML + c * (cwid + gap)
        y = top + r * (height + gap)
        rect(slide, x, y, cwid, height, fill=fill, radius=radius, line=line)
        if accent_bar:
            rect(slide, x, y, 0.055, height, fill=BLUE, radius=None)
        ix = x + pad + (0.06 if accent_bar else 0)
        iw = cwid - 2 * pad - (0.06 if accent_bar else 0)
        # title height from its own estimated wrap, so a long title in a narrow card
        # pushes the body down instead of colliding with it
        tlines = max(1, math.ceil(len(t) / max(8, (iw * 72.0) / (tsize * 0.5)))) if t else 0
        th = tlines * tsize * 1.16 / 72.0
        if t:
            txt(slide, ix, y + pad - 0.02, iw, th, t, size=tsize, color=tcolor,
                bold=True, line_spacing=1.14)
        by = y + pad - 0.02 + (th + 0.07 if t else 0)
        bh = height - (by - y) - pad
        if isinstance(b, list):
            paras = [[("·  ", {"color": BLUE, "bold": True}), (s, {})] for s in b]
            txt(slide, ix, by, iw, bh, paras, size=bsize, color=bcolor,
                line_spacing=1.24, space_after=3.2)
        elif b:
            txt(slide, ix, by, iw, bh, b, size=bsize, color=bcolor, line_spacing=1.3)


def stat_tiles(slide, items, top, height=1.02, cols=None, num_size=27, lab_size=8.2,
               dark=False):
    """items: list of (number, label)."""
    cols = cols or len(items)
    gap = 0.18
    cwid = (CW - gap * (cols - 1)) / cols
    for i, (num, lab) in enumerate(items):
        x = ML + i * (cwid + gap)
        rect(slide, x, top, cwid, height,
             fill=None if dark else G50, radius=20,
             line=BLUE_800 if dark else None)
        txt(slide, x + 0.16, top + 0.13, cwid - 0.32, 0.44, num, size=num_size,
            color=BLUE if not dark else WHITE, bold=True, align=PP_ALIGN.CENTER,
            line_spacing=1.0)
        txt(slide, x + 0.12, top + 0.62, cwid - 0.24, height - 0.72, lab,
            size=lab_size, color=G600 if not dark else BLUE_200,
            align=PP_ALIGN.CENTER, line_spacing=1.22)


def table(slide, headers, rows, top, widths=None, height=None, hsize=8.6, bsize=8.4,
          row_h=None, bold_first=False):
    ncol = len(headers)
    widths = widths or [1.0 / ncol] * ncol
    tot = sum(widths)
    widths = [w / tot * CW for w in widths]
    hh = 0.30
    row_h = row_h or min(0.335, (SH - 0.66 - top - hh) / max(1, len(rows)))
    # header
    x = ML
    rect(slide, ML, top, CW, hh, fill=NAVY, radius=None)
    for i, hd in enumerate(headers):
        txt(slide, x + 0.11, top + 0.055, widths[i] - 0.2, hh - 0.08, hd, size=hsize,
            color=WHITE, bold=True, fit_check=False)
        x += widths[i]
    # body
    y = top + hh
    for j, row in enumerate(rows):
        rect(slide, ML, y, CW, row_h, fill=G50 if j % 2 == 0 else WHITE, radius=None)
        rect(slide, ML, y, CW, 0.008, fill=G200, radius=None)
        x = ML
        for i, cell in enumerate(row):
            b = bold_first and i == 0
            txt(slide, x + 0.11, y + 0.045, widths[i] - 0.2, row_h - 0.06, cell,
                size=bsize, color=NAVY if b else G700, bold=b, line_spacing=1.16,
                fit_check=False)
            x += widths[i]
        y += row_h
    rect(slide, ML, y, CW, 0.008, fill=G200, radius=None)
    return y


def two_col(slide, left, right, top, height=None, lfill=G50, rfill=None,
            lline=None, rline=None, tsize=12, bsize=9.0):
    """left/right: (title, [bullets], title_color)"""
    gap = 0.24
    w = (CW - gap) / 2
    height = height or (SH - 0.66 - top)
    for k, (blk, fill, ln) in enumerate(((left, lfill, lline), (right, rfill, rline))):
        t, items, tc = blk
        x = ML + k * (w + gap)
        rect(slide, x, top, w, height, fill=fill, radius=20, line=ln)
        txt(slide, x + 0.22, top + 0.19, w - 0.44, 0.28, t, size=tsize, color=tc,
            bold=True)
        paras = [[("·  ", {"color": tc, "bold": True}), (s, {})] for s in items]
        txt(slide, x + 0.22, top + 0.58, w - 0.44, height - 0.78, paras, size=bsize,
            color=G700, line_spacing=1.26, space_after=3.4)


def process(slide, steps, top, height=1.5, tsize=10.5, bsize=8.4, numbered=True):
    """steps: list of (title, body). Chevron flow with arrows between."""
    n = len(steps)
    arrow = 0.26
    cwid = (CW - arrow * (n - 1)) / n
    for i, (t, b) in enumerate(steps):
        x = ML + i * (cwid + arrow)
        rect(slide, x, top, cwid, height, fill=G50, radius=20)
        rect(slide, x, top, cwid, 0.05, fill=BLUE, radius=None)
        if numbered:
            cir = rect(slide, x + 0.16, top + 0.19, 0.30, 0.30, fill=BLUE, radius=9999)
            try:
                cir.adjustments[0] = 0.5
            except Exception:
                pass
            txt(slide, x + 0.16, top + 0.245, 0.30, 0.22, str(i + 1), size=10,
                color=WHITE, bold=True, align=PP_ALIGN.CENTER, fit_check=False)
            ty = top + 0.19
            tx = x + 0.54
            tw = cwid - 0.70
        else:
            ty, tx, tw = top + 0.19, x + 0.18, cwid - 0.36
        txt(slide, tx, ty, tw, 0.30, t, size=tsize, color=NAVY, bold=True,
            line_spacing=1.12)
        txt(slide, x + 0.18, top + 0.62, cwid - 0.36, height - 0.80, b, size=bsize,
            color=G700, line_spacing=1.26)
        if i < n - 1:
            txt(slide, x + cwid, top + height / 2 - 0.14, arrow, 0.28, "›",
                size=17, color=BLUE_300, bold=True, align=PP_ALIGN.CENTER,
                fit_check=False)


def banner(slide, label, text, top, height=0.72, fill=None, dark=True, tsize=14.5):
    """A navy or blue statement band."""
    b = rect(slide, ML, top, CW, height, fill=fill or NAVY, radius=20)
    if dark and fill is None:
        set_radial_gradient(b, [(0, "1B36B8"), (55000, "08186B"), (100000, "050F40")],
                            center=(28000, 50000))
    if label:
        txt(slide, ML + 0.24, top + 0.13, 1.9, 0.2, label, size=8, color=BLUE_300,
            bold=True, fit_check=False)
        txt(slide, ML + 0.24, top + 0.34, CW - 0.5, height - 0.46, text, size=tsize,
            color=WHITE, bold=True, line_spacing=1.14)
    else:
        txt(slide, ML + 0.24, top + 0.1, CW - 0.5, height - 0.2, text, size=tsize,
            color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.16)


def divider(slide, num, title, sub):
    navy_bg(slide, center=(30000, 40000))
    txt(slide, ML, 2.5, CW, 0.9, num, size=64, color=BLUE_400, bold=True,
        line_spacing=1.0, fit_check=False)
    txt(slide, ML, 3.42, CW, 0.66, title, size=36, color=WHITE, bold=True,
        line_spacing=1.06)
    txt(slide, ML, 4.24, CW - 3.2, 0.6, sub, size=13, color=BLUE_200,
        line_spacing=1.3)
    rect(slide, ML, 4.06, 1.5, 0.045, fill=BLUE, radius=None)


# ----------------------------------------------------------------------------- deck
prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)


def set_theme_font(presentation, name):
    """Point the theme's major and minor latin fonts at Satoshi, so the declared
    typeface survives the trip through Google Slides even where it substitutes."""
    part = presentation.slide_masters[0].part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
    )
    theme = etree.fromstring(part.blob)
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for scheme in theme.iter(f"{ns}fontScheme"):
        for which in ("majorFont", "minorFont"):
            for f in scheme.iter(f"{ns}{which}"):
                for latin in f.iter(f"{ns}latin"):
                    latin.set("typeface", name)
    part._blob = etree.tostring(theme, xml_declaration=True, encoding="UTF-8",
                                standalone=True)
    part.blob = part._blob


set_theme_font(prs, FONT)
dark_slides = set()
S = []  # builders, so page numbers can be added after totals are known


def page(fn, dark=False):
    S.append((fn, dark))


# ---- 1 cover
def s_cover(sl):
    navy_bg(sl, center=(62000, 18000))
    txt(sl, ML, 2.28, CW, 0.3, "FITXPRESS BY 3DLOOK", size=11, color=BLUE_300,
        bold=True)
    txt(sl, ML, 2.70, 9.6, 1.5, "Product Marketing\nPlaybook", size=52, color=WHITE,
        bold=True, line_spacing=1.02)
    rect(sl, ML, 4.42, 1.7, 0.05, fill=BLUE, radius=None)
    txt(sl, ML, 4.70, 9.0, 0.3, "Nika  ·  Product Marketing Manager  ·  Onboarding",
        size=15, color=BLUE_100)
    txt(sl, ML, 5.08, 9.6, 0.5,
        "Presentations and landing pages, end to end, for every FitXpress health use case.",
        size=12, color=BLUE_200, line_spacing=1.3)
    txt(sl, ML, 6.42, 7.0, 0.24,
        "August 2026  ·  Version 1.0  ·  CONFIDENTIAL, internal use only",
        size=9, color=BLUE_300, bold=True)
    txt(sl, SW - ML - 3.0, 6.42, 3.0, 0.24, "Structured after the FitXpress Sales Playbook",
        size=9, color=BLUE_300, align=PP_ALIGN.RIGHT)


page(s_cover, dark=True)


# ---- 2 what's inside
def s_toc(sl):
    slide_head(sl, "NAVIGATION", "What's Inside",
               "Nine sections. Section 4 is the working reference. Sections 5 and 6 are the two production lines.")
    items = [
        ("01", "Role and scope"),
        ("02", "Brand and voice guardrails"),
        ("03", "Design system quick reference"),
        ("04", "Use-case library, 11 health segments"),
        ("05", "Deck build playbook"),
        ("06", "Landing-page build playbook"),
        ("07", "Messaging and proof library"),
        ("08", "QA checklist"),
        ("09", "Open items to resolve"),
    ]
    gap, cols = 0.2, 3
    w = (CW - gap * 2) / 3
    h = 1.28
    for i, (n, t) in enumerate(items):
        r, c = divmod(i, cols)
        x = ML + c * (w + gap)
        y = 2.12 + r * (h + gap)
        rect(sl, x, y, w, h, fill=G50, radius=20)
        rect(sl, x, y, 0.055, h, fill=BLUE, radius=None)
        txt(sl, x + 0.26, y + 0.22, 1.0, 0.4, n, size=23, color=BLUE_200, bold=True,
            fit_check=False)
        txt(sl, x + 0.26, y + 0.68, w - 0.5, 0.44, t, size=12.5, color=NAVY, bold=True,
            line_spacing=1.16)


page(s_toc)


# ---- 3 the one rule
def s_rule(sl):
    navy_bg(sl, center=(72000, 30000))
    eyebrow(sl, "THE RULE ABOVE ALL THE OTHERS", color=BLUE_300)
    txt(sl, ML, 0.86, 8.6, 1.5,
        "Every figure in a deck or on a page\ncomes from proof-points.md.",
        size=31, color=WHITE, bold=True, line_spacing=1.1)
    txt(sl, ML, 2.42, 8.4, 0.72,
        "A number that is not there does not ship. It goes into an Open items block, and the slot "
        "stays visibly empty until Vadim confirms it. This is the fastest way this function either "
        "earns or loses credibility with a buyer doing diligence.",
        size=12, color=BLUE_100, line_spacing=1.36)
    steps = [
        ("Need a number", "A deck slot or a page section calls for a figure"),
        ("Check proof-points.md", "The approved proof set in section 7 is the whole permitted universe"),
        ("Not there?", "Do not approximate, do not borrow from an adjacent vertical, do not average"),
        ("Open items block", "Record it, leave the slot empty, route it to Vadim"),
    ]
    n = len(steps)
    arrow, top, h = 0.26, 3.52, 1.32
    w = (CW - arrow * (n - 1)) / n
    for i, (t, b) in enumerate(steps):
        x = ML + i * (w + arrow)
        c = rect(sl, x, top, w, h, fill=None, radius=20, line=BLUE_700, line_w=1.1)
        rect(sl, x, top, w, 0.05, fill=BLUE, radius=None)
        txt(sl, x + 0.2, top + 0.22, w - 0.4, 0.3, t, size=11.5, color=WHITE, bold=True)
        txt(sl, x + 0.2, top + 0.6, w - 0.4, h - 0.78, b, size=8.6, color=BLUE_200,
            line_spacing=1.3)
        if i < n - 1:
            txt(sl, x + w, top + h / 2 - 0.14, arrow, 0.28, "›", size=17,
                color=BLUE_300, bold=True, align=PP_ALIGN.CENTER, fit_check=False)
    txt(sl, ML, 5.30, CW, 0.3,
        "The same discipline applies to customer names and quotes: they come from case-studies/, "
        "and only with current written permission.",
        size=10, color=BLUE_300, italic=True)


page(s_rule, dark=True)


# ============================================================ 01 ROLE
def s_div1(sl):
    divider(sl, "01", "Role and scope",
            "What Nika owns, who she needs, where the assets live, and what the first two weeks look like.")


page(s_div1, dark=True)


def s_owns(sl):
    top = slide_head(sl, "01 · ROLE AND SCOPE", "What Nika owns, and what routes elsewhere",
                     "Owning the wrong thing is as expensive as owning nothing. Everything below the line has another owner or another pipeline.")
    two_col(
        sl,
        ("Owned end to end", [
            "Vertical demo decks, FitXpress for X. Copy, structure, slide-level design brief, versioning",
            "First-call and leave-behind decks, short cuts of the demo deck on the same spine",
            "Vertical landing pages on 3dlook.ai, as the driver through the /page pipeline",
            "Campaign landing pages, same section map, shorter, one message each",
            "Two-pagers and one-pagers, using the print variant of the design system",
            "Messaging per segment: hero line, proof set, objection answers, disclaimer",
        ], BLUE),
        ("Routes elsewhere", [
            "Blog and SEO articles, route to /new-article. A page brief and an article brief resolve against different plans and will contradict each other",
            "Social posts from a published article, route to /post-from-article",
            "Outbound sequences, route to /outbound",
            "Pricing decisions and exceptions: Katerina Galich or Kate Kondakova",
            "Final editorial sign-off on published health copy: Asselya Sekerova, non-optional",
            "Claims touching medical, compliance, underwriting or employment: escalate to Whitney Cathcart and Katerina before publishing",
        ], NAVY),
        top, lfill=BLUE_50, rfill=G50, bsize=8.6,
    )


page(s_owns)


def s_team(sl):
    top = slide_head(sl, "01 · ROLE AND SCOPE", "Who Nika needs, and what for",
                     "Weekly: what shipped, what is blocked, which numbers need confirmation. Bi-weekly with Katerina on positioning. Monthly asset audit.")
    cards(sl, [
        ("Katerina Galich · CEO",
         "Positioning calls, competitive framing sign-off, pricing exceptions, customer logo and quote permission"),
        ("Whitney Cathcart · Co-founder and CCO",
         "Deck guidance, and claims framing on every regulated vertical"),
        ("Vadim Bilan · Marketing",
         "Number confirmation, proof-points.md updates, brand asset access, the repo and the automation pipelines"),
        ("Asselya Sekerova · Editorial",
         "Final review of health copy before publish. Non-optional"),
        ("Max Kucherenko · VP Product",
         "Capability questions, roadmap status, what may be promised"),
        ("Nick, Olena, Kateryna · Business development",
         "The real objections and the vertical specifics a page or deck needs. The source for what a buyer in this market actually asks"),
    ], top, cols=3, height=1.42, accent_bar=True, tsize=11)
    banner(sl, "ESCALATE IMMEDIATELY",
           "A deck or page needs a figure that does not exist  ·  a customer name or quote is requested "
           "without written permission  ·  a claim starts drifting toward diagnosis, eligibility "
           "decisioning or guaranteed compliance",
           top + 3.20, height=0.90, tsize=11.5)


page(s_team)


def s_assets(sl):
    top = slide_head(sl, "01 · ROLE AND SCOPE", "Where the assets live",
                     "Local repo governs voice, claims and design. Drive holds the existing collateral and the raw buyer material.")
    txt(sl, ML, top, CW, 0.24, "LOCAL REPO, THE SOURCE OF TRUTH", size=8.6, color=BLUE,
        bold=True)
    table(sl, ["Path", "What it governs"], [
        ["proof-points.md", "Every number. Nothing else is a citable figure"],
        ["about-me.md", "Brand voice, claims discipline, accuracy framing, call-to-action discipline"],
        ["audience.md", "The seven health segments: the hook, and what NOT to say"],
        ["DESIGN.md", "Every design token. The only design source of truth"],
        ["CLAUDE.md", "Company facts, the two-product split, banned words, compliance posture"],
        ["icp-detail.md · messaging.md", "Buyer titles, revenue thresholds, triggers. Approved hero lines"],
        ["case-studies/", "The only place a customer name and a customer number may come from"],
        ["editorial-guardrails.md · ai-tells-sweep.md", "The 11 principles, rules M1 and M2, the hard-fail list and the detector"],
        ["page-builder/ · quality-rubric.md", "The page pipeline gates and the 17-slot kit. The 20-point QC rubric"],
    ], top + 0.30, widths=[0.30, 0.70], row_h=0.272)
    txt(sl, ML, top + 3.06, CW, 0.24, "DRIVE SALES FOLDER, THE SAME ONE SALES USES",
        size=8.6, color=BLUE, bold=True)
    cards(sl, [
        ("Marketing docs & Decks",
         "Insurance underwriting and telehealth demo decks (May 2026), the 2-pager, integration guidelines"),
        ("Technology & Accuracy",
         "Accuracy and Repeatability Analysis (April 2026), fat estimation approach, earlier studies"),
        ("Useful info",
         "Body Mass Index (BMI) verification use case, Prism Labs analysis, the market-sizing Brainstorm deck"),
        ("Nick Omelchak, call transcripts",
         "Real buyer objections in the buyer's own words. The best raw material for a problem slide or a vertical FAQ"),
    ], top + 3.34, cols=4, height=1.14, tsize=9.8, bsize=8.0, fill=BLUE_50)
    txt(sl, ML, top + 4.58, CW, 0.40,
        "Two cautions. FX Guidance for Decks (Whitney, February 2026) is a shortcut whose target did not "
        "resolve, so ask for a re-share. The published 2-pager says \"seamlessly embeds\" and \"under 60 "
        "seconds\", both now wrong, so treat it as a layout reference and not as copy.",
        size=8.4, color=G600, italic=True, line_spacing=1.28)


page(s_assets)


def s_weeks(sl):
    top = slide_head(sl, "01 · ROLE AND SCOPE", "The first two weeks",
                     "Read, then audit, then ship one real thing. Seeing every gate work once is worth more than reading about all of them.")
    process(sl, [
        ("Days 1 to 3 · Read and calibrate",
         "CLAUDE.md, about-me.md, audience.md, DESIGN.md, proof-points.md, editorial-guardrails.md. "
         "Then the insurance underwriting deck copy end to end as the structural exemplar, and the "
         "telehealth page structure document as the landing-page exemplar. Then the sales playbook, "
         "to hear how sales speaks."),
        ("Days 4 to 7 · Audit",
         "Score the existing decks and the six live vertical pages against section 8. Produce one list: "
         "every figure in circulation that is not in proof-points.md, every banned word on a live page, "
         "every claim that needs Whitney. Section 9 is the starting point for that list, not the end of it."),
        ("Week 2 · Ship one thing",
         "Rewrite /for-bmi-verification/, the weakest live page and the vertical with a real customer "
         "reference in UK Meds, taking it through the full pipeline in section 6. It is the shortest "
         "path to seeing every gate work."),
    ], top + 0.1, height=2.32, tsize=12.5, bsize=9.6)
    banner(sl, None,
           "The week-two page rewrite is deliberately the hardest live page, not the easiest. "
           "A gate that has never blocked anything has not been tested.",
           top + 2.76, height=0.78, tsize=12.5)


page(s_weeks)


# ============================================================ 02 VOICE
def s_div2(sl):
    divider(sl, "02", "Brand and voice guardrails",
            "Positioning, the voice, claims discipline, the accuracy framing, banned language, and the disclaimer set.")


page(s_div2, dark=True)


def s_pos(sl):
    top = slide_head(sl, "02 · BRAND AND VOICE", "Outcomes and workflow, not accuracy")
    banner(sl, "THE POSITIONING, IN ONE LINE",
           "FitXpress is a trusted workflow layer for verified body data, for industries where "
           "measurement consistency, fraud reduction and audit-ready records drive a business result. "
           "It is not a measurement API.",
           top - 0.34, height=1.16, tsize=14)
    two_col(
        sl,
        ("Why this matters commercially", [
            "Body scanning commoditizes on a 12 to 36 month horizon as vision foundation models become common",
            "Apple or Google may ship native body-measurement primitives",
            "A deck that wins on \"most accurate\" loses the day that happens",
            "A deck that wins on workflow integration, governance, auditability and longitudinal reliability does not",
        ], NAVY),
        ("The consequence for every artifact", [
            "The hero line names the buyer's business outcome",
            "Accuracy appears later, scoped, as support for that outcome",
            "Never lead with \"most accurate body scanning\", \"AI-powered body scanning\" alone, \"replace your in-person fitting\", or \"just plug in our API\"",
            "Prefer \"verified body data\" over \"accurate measurements\", and \"trusted workflow layer\" over \"scanning API\"",
        ], BLUE),
        top + 1.02, height=2.42, lfill=G50, rfill=BLUE_50,
    )
    txt(sl, ML, top + 3.60, CW, 0.40,
        "Master positioning, for the About slide and the page footer: 3DLOOK is the trusted "
        "infrastructure for verified body data, built for industries where measurement consistency, "
        "fraud reduction and audit-ready workflows drive business outcomes.",
        size=9.4, color=G600, italic=True, line_spacing=1.3)


page(s_pos)


def s_voice(sl):
    top = slide_head(sl, "02 · BRAND AND VOICE", "The voice",
                     "Calm, specific, evidence-led B2B. It sells by clarifying the buyer's decision, and it names its own limits before anyone else can.")
    cards(sl, [
        ("The reframe move, the signature",
         "Open by turning the obvious question into the better one. \"How accurate is it?\" becomes "
         "\"accurate enough for which decision?\". \"DEXA or mobile?\" becomes \"how do the two fit "
         "together inside one program?\""),
        ("Declarative and unhurried",
         "Mostly 15 to 30 word sentences, 2 to 4 sentence paragraphs. An occasional short verdict line "
         "for emphasis: \"In short.\" \"Production conditions are not lab conditions.\""),
        ("Concrete over abstract",
         "Every claim carries a number, a named source, a condition or a disclosed limitation. A vague "
         "adjective is missing information."),
        ("Honest about limits",
         "State what the product does not do in the same breath as what it does. This is what makes the "
         "rest of the claims credible."),
        ("Buyer framing",
         "\"Enterprise teams\", \"insurers\", \"care teams\", \"programs\". Use \"you\" sparingly, mainly "
         "on landing pages and conversion sections addressing a decision-maker directly."),
        ("Neutral authority",
         "Cite external bodies (CDC, Munich Re, Swiss Re, NAIC, LIMRA, ISO) with links, rather than "
         "asserting authority in our own voice."),
        ("Compare by role, not by hype",
         "A comparison answers \"which method fits which workflow?\". A clean sweep across every row "
         "reads as marketing, and buyers discount it."),
        ("No jokes in published copy",
         "Internal strategy documents may be witty. Customer-facing copy is sober."),
    ], top, cols=4, height=1.62, accent_bar=True, tsize=10.5, bsize=8.3)


page(s_voice)


def s_claims(sl):
    top = slide_head(sl, "02 · BRAND AND VOICE", "Claims discipline: what FitXpress is, and is not",
                     "The left column is a hard fail in any artifact. The right column is the approved substitute, and it is not weaker copy.")
    table(sl, ["Never claim", "Say instead"], [
        ["Diagnose a condition", "Provides structured body data that supports clinician review"],
        ["Make treatment, underwriting, hiring or fitness-for-duty decisions",
         "Supporting evidence, not standalone decisioning"],
        ["Replace clinicians, dual-energy X-ray absorptiometry (DEXA), bioelectrical impedance analysis (BIA), calibrated scales or a protocol reference method",
         "Complements the reference method by collecting remote measurements between clinical assessment points"],
        ["Guarantee regulatory compliance", "Supports compliant workflows and audit-ready documentation"],
        ["Detect fraud automatically", "Supports fraud-prevention workflows by flagging discrepancies for review"],
        ["Act as a standalone medical authority", "An intake and documentation layer, alongside clinical oversight"],
    ], top, widths=[0.46, 0.54], row_h=0.44, bold_first=True)
    two_col(
        sl,
        ("The medical framing, exactly", [
            "\"FitXpress is not positioned as a medical device.\" In that wording",
            "Never \"medical device regulation does not apply\"",
            "Compliance is framed on data-privacy frameworks (HIPAA, GDPR, SOC 2 where applicable), not on medical-device frameworks (FDA Class II, CE-MDR)",
        ], NAVY),
        ("Never claimed at all", [
            "Food and Drug Administration (FDA) clearance",
            "Medical advice or diagnosis",
            "SOC 2 certification, which is in progress. Confirm with Vadim before any mention",
            "Direct processing of protected health information",
        ], BLUE),
        top + 2.98, height=1.62, lfill=G50, rfill=BLUE_50, bsize=8.4, tsize=11,
    )


page(s_claims)


def s_repeat(sl):
    top = slide_head(sl, "02 · BRAND AND VOICE", "Repeatability is not accuracy",
                     "The single framing most often broken, and the one that sells. Mixing the two produces a figure that describes nothing.")
    two_col(
        sl,
        ("Accuracy", [
            "Measured against a reference method",
            "Only meaningful with four conditions named: reference method, measurement protocol, population tested, intended workflow",
            "Never reduced to one universal number. Always qualify: accurate for which decision, against which reference, under which capture protocol, for which population, at what tolerance",
        ], NAVY),
        ("Repeatability", [
            "Scan-to-scan consistency",
            "What matters for longitudinal use: glucagon-like peptide-1 (GLP-1) progress over 30, 60 and 90 days, year-over-year underwriting refreshes, survivorship monitoring",
            "Written as \"< 1 cm\". Locked convention",
        ], BLUE),
        top, height=1.72, lfill=G50, rfill=BLUE_50, bsize=8.6, tsize=12,
    )
    txt(sl, ML, top + 1.88, CW, 0.24, "TWO BENCHMARKS, NEVER COMBINED", size=8.6,
        color=BLUE, bold=True)
    b1 = rect(sl, ML, top + 2.16, (CW - 0.24) / 2, 1.10, fill=WHITE, radius=20,
              line=BLUE_200, line_w=1.2)
    txt(sl, ML + 0.22, top + 2.32, (CW - 0.24) / 2 - 0.44, 0.24,
        "1 · Internal validation, against expert manual measurement", size=10,
        color=NAVY, bold=True)
    txt(sl, ML + 0.22, top + 2.62, (CW - 0.24) / 2 - 0.44, 0.56,
        "96 to 97% accuracy, typical absolute error 1.5 to 2.0 cm, scan-to-scan repeatability "
        "< 1 cm. USABLE: this set is in proof-points.md.",
        size=8.8, color=G700, line_spacing=1.28)
    x2 = ML + (CW - 0.24) / 2 + 0.24
    rect(sl, x2, top + 2.16, (CW - 0.24) / 2, 1.10, fill=G100, radius=20,
         line=G300, line_w=1.2)
    txt(sl, x2 + 0.22, top + 2.32, (CW - 0.24) / 2 - 0.44, 0.24,
        "2 · The ISO multi-company benchmark", size=10, color=G700, bold=True)
    txt(sl, x2 + 0.22, top + 2.62, (CW - 0.24) / 2 - 0.44, 0.56,
        "Reported as 0.40 cm session-to-session repeatability. NOT in proof-points.md, though it is "
        "in both shipped decks. Treat as pending. See section 9, item 2.",
        size=8.8, color=G700, line_spacing=1.28)
    banner(sl, "THE REUSABLE PARAGRAPH, ADAPT THE VERTICAL AND THE DECISION",
           "The better diligence question is: accurate enough for which decision? For [decision in this "
           "vertical], what matters is whether a change between scans is real rather than measurement "
           "noise. Against expert manual measurement, FitXpress reaches 96 to 97% accuracy with a typical "
           "absolute error of 1.5 to 2.0 cm, and scan-to-scan variance stays < 1 cm. Detailed methodology "
           "is available under NDA. FitXpress is not positioned as a medical device, and it supports "
           "clinician review rather than replacing it.",
           top + 3.42, height=1.20, tsize=10.5)


page(s_repeat)


def s_words(sl):
    top = slide_head(sl, "02 · BRAND AND VOICE", "Words to use, words never to use",
                     "The banned list is a grep, so it is cheap to check and there is no excuse for a miss. The constructions are the expensive part.")
    two_col(
        sl,
        ("Use", [
            "Operational verbs: supports, helps standardize, provides structured records, reduces manual intake, standardizes capture, supports review, improves documentation consistency, reduces rework, supports scan-to-scan comparison",
            "Precise hedges: designed to, can support, where the workflow or protocol allows, supporting evidence rather than standalone decisioning, not a replacement for clinician review, an intake and documentation layer",
            "Framing phrases: accurate enough for which decision? · compare by role, not by hype",
            "\"Supports clinician review\" is the workhorse phrase. Use it often, and honestly",
        ], BLUE),
        ("Never use", [
            "Banned words: leverage, utilize, harness, robust, seamless, comprehensive, revolutionize, cutting-edge, state-of-the-art, game-changer, disrupt, delve, tapestry, realm, groundbreaking, best-in-class, industry-leading, world-class, unparalleled, and figurative \"navigate\"",
            "Banned phrasings: \"In today's fast-paced world\", \"Unlock the power of\", \"Are you struggling with?\", \"It's no secret that\", \"Let's dive in\", and \"AI-powered\" standing alone as a value claim",
            "Asselya's word rules: never \"objective\" about our own output, never \"the reader\" or \"this guide\", never \"by hand\", never \"plus\" as a connector, never \"let\", never \"so\" introducing a benefit",
        ], NAVY),
        top, height=2.10, lfill=BLUE_50, rfill=G50, bsize=8.4, tsize=12,
    )
    txt(sl, ML, top + 2.26, CW, 0.24, "BANNED CONSTRUCTIONS, AND THE FIX", size=8.6,
        color=BLUE, bold=True)
    table(sl, ["Construction", "Why", "Fix"], [
        ["\"It's not just X, it's Y\", \"not only X but also Y\"",
         "The most reliable AI signature in English marketing copy, and it survives casual editing",
         "State Y directly"],
        ["Adjectival punch triads: \"fast, reliable, scalable\"",
         "Reads generated. A list of real things (\"positioning, posture, equipment\") is a list and is fine",
         "Two adjectives, or one with a condition"],
        ["Em dash and en dash",
         "Banned in all contexts, no exceptions",
         "Comma, full stop, or brackets"],
        ["Corrective negation \"X, not Y\"",
         "Sounds corrective and dismissive. Exception: a necessary product, clinical, legal or regulatory boundary",
         "Lead with the recommended approach and explain its purpose"],
    ], top + 2.54, widths=[0.28, 0.46, 0.26], row_h=0.40)
    txt(sl, ML, top + 4.30, CW, 0.36,
        "Rules M1 and M2. M1: expand every abbreviation at first use, including the obvious ones and "
        "every cited regulator. M2: one clear negative statement of scope, stated once, with no second "
        "negation chained onto it. The AI-tells sweep is a separate pass after writing, never during.",
        size=8.6, color=G600, italic=True, line_spacing=1.28)


page(s_words)


def s_disclaim(sl):
    top = slide_head(sl, "02 · BRAND AND VOICE", "Disclaimer boilerplate for sensitive verticals",
                     "A scope note early, not a disclaimer bolted on at the end, plus an italic line near any accuracy or eligibility claim. Keep the first sentence intact, adapt the second clause.")
    banner(sl, "UNIVERSAL BASE, EVERY SENSITIVE VERTICAL",
           "FitXpress is not positioned as a medical device. It provides structured body data that "
           "supports clinical and operational workflows.",
           top - 0.30, height=0.86, tsize=13)
    table(sl, ["Vertical", "The second clause"], [
        ["Telehealth, GLP-1 and weight loss",
         "It supports remote progress tracking and clinician review; treatment decisions remain with the care team"],
        ["Online pharmacy and remote prescribing",
         "It supports eligibility verification workflows and audit-ready documentation; prescribing and eligibility determinations remain with the prescriber"],
        ["Insurance underwriting",
         "Outputs are supporting evidence for underwriter review; risk classification and pricing decisions remain with the carrier"],
        ["Health plans and employer wellness",
         "It supports remote verification and documentation for incentive programs; reward and eligibility determinations remain with the program administrator"],
        ["Bariatric and metabolic clinics",
         "It supports remote pre-qualification workflows; surgical candidacy remains the clinical team's determination"],
        ["Occupational health",
         "It supports screening intake and documentation; fitness-for-duty and clearance determinations remain with the occupational health provider"],
        ["Clinical trials",
         "It standardizes and documents anthropometric capture; endpoint validation and protocol compliance remain with the sponsor and the investigator"],
        ["Plastic surgery",
         "It supports remote pre-screening and planning input; surgical and anaesthetic risk assessment remains with the surgeon"],
        ["BCRL and oncology survivorship",
         "It supports remote monitoring workflows and reproducible body records; it does not detect or diagnose lymphedema, and clinical assessment remains with the care team"],
    ], top + 0.72, widths=[0.26, 0.74], row_h=0.325, bold_first=True)
    txt(sl, ML, top + 3.86, CW, 0.44,
        "Always safe: \"supports compliant workflows\", \"audit-ready documentation\", \"HIPAA-aware and "
        "GDPR-aligned\". Never \"makes you compliant\", never \"guarantees compliance\". Every control is "
        "stated with its limit: pose validation, clothing detection and live capture reduce error and "
        "fraud exposure, and they do not remove the need for capture instructions, retake logic or "
        "deployment-specific thresholds. Breast cancer-related lymphedema (BCRL) is the most "
        "compliance-sensitive vertical: lead with the scope note, and route every draft to Whitney.",
        size=8.4, color=G600, italic=True, line_spacing=1.3)


page(s_disclaim)


# ============================================================ 03 DESIGN
def s_div3(sl):
    divider(sl, "03", "Design system quick reference",
            "DESIGN.md is the only design source of truth. Use tokens exactly as written. Never introduce a font, colour or radius that is not in it.")


page(s_div3, dark=True)


def s_colour(sl):
    top = slide_head(sl, "03 · DESIGN SYSTEM", "Colour, and how it is weighted",
                     "Clinical precision with consumer-app polish. Restraint over density, product over icons, precision as a visual language, depth rather than flat fills.")
    core = [("#143DFF", "Electric blue\nThe single sharp accent", BLUE, WHITE),
            ("#050F40", "Navy\nHero, proof bands, CTA band, footer", NAVY, WHITE),
            ("#000000", "Black\nDark buttons, high-contrast type", BLACK, WHITE),
            ("#FFFFFF", "White\nDominant on content zones", WHITE, NAVY)]
    w = (CW - 0.18 * 3) / 4
    for i, (hexv, lab, fill, tc) in enumerate(core):
        x = ML + i * (w + 0.18)
        rect(sl, x, top, w, 1.16, fill=fill, radius=20,
             line=G300 if fill == WHITE else None, line_w=1.0)
        txt(sl, x + 0.2, top + 0.20, w - 0.4, 0.3, hexv, size=15, color=tc, bold=True,
            fit_check=False)
        txt(sl, x + 0.2, top + 0.58, w - 0.4, 0.5, lab, size=8.4,
            color=tc if fill != WHITE else G600, line_spacing=1.26)
    # scales
    y = top + 1.36
    txt(sl, ML, y, CW, 0.22, "BLUE SCALE, TEN STEPS", size=8, color=BLUE, bold=True)
    blues = ["ECEFFF", "D8DEFF", "B1BDFF", "8A9CFF", "4F6DFF", "143DFF", "0F2ECD",
             "0B2299", "08186B", "050F40"]
    sw = (CW - 0.06 * 9) / 10
    for i, h in enumerate(blues):
        c = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        rect(sl, ML + i * (sw + 0.06), y + 0.26, sw, 0.40, fill=c, radius=15,
             line=G300 if i == 0 else None)
        txt(sl, ML + i * (sw + 0.06), y + 0.70, sw, 0.18, "#" + h, size=6.2,
            color=G500, align=PP_ALIGN.CENTER, fit_check=False)
    y2 = y + 1.00
    txt(sl, ML, y2, CW, 0.22, "GRAY SCALE, TEN STEPS", size=8, color=BLUE, bold=True)
    grays = ["F9F9F9", "F2F2F2", "E5E5E5", "D1D1D1", "A8A8A8", "808080", "666666",
             "4C4C4C", "333333", "1A1A1A"]
    for i, h in enumerate(grays):
        c = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        rect(sl, ML + i * (sw + 0.06), y2 + 0.26, sw, 0.40, fill=c, radius=15,
             line=G300 if i < 2 else None)
        txt(sl, ML + i * (sw + 0.06), y2 + 0.70, sw, 0.18, "#" + h, size=6.2,
            color=G500, align=PP_ALIGN.CENTER, fit_check=False)
    cards(sl, [
        ("Weighting",
         "White dominates content zones. Navy #050F40 carries 60 to 70% of the weight on hero, proof, "
         "CTA and footer. Electric blue stays a single sharp accent and never becomes a large fill."),
        ("Navy glow gradient",
         "Radial glow for hero, CTA and footer: brighter blue upper-centre falling to deep navy at the "
         "edges. Stops #4F6DFF or #0B2299 core, then #08186B, then #050F40. Add subtle grain or a faint "
         "measurement-grid texture."),
        ("On #2962FF",
         "Superseded. It was a placeholder before the Figma export arrived, and both CLAUDE.md and "
         "DESIGN.md retire it along with Inter. The canonical accent is #143DFF and the canonical "
         "typeface is Satoshi. Status colours are not brand tokens: if a status UI is needed, choose an "
         "accessible green or amber and flag it."),
    ], y2 + 1.00, cols=3, height=1.48, tsize=10.5, bsize=8.2, accent_bar=True)


page(s_colour)


def s_type(sl):
    top = slide_head(sl, "03 · DESIGN SYSTEM", "Typography, spacing, radius, buttons",
                     "Satoshi for headings and body. Weights 400, 500, 600, 700, and 900 for display accents.")
    # type scale table
    txt(sl, ML, top, CW * 0.48, 0.22, "TYPE SCALE", size=8, color=BLUE, bold=True)
    tw = CW * 0.48
    rows = [["Display", "80px", "700"], ["H1", "65px", "700"], ["H2", "50px", "700"],
            ["H3", "40px", "600"], ["H4", "35px", "600"], ["H5", "27px", "500"],
            ["Body lg", "20px", "400"], ["Body", "17px", "400"], ["Body sm", "12px", "400"]]
    y = top + 0.28
    rect(sl, ML, y, tw, 0.26, fill=NAVY)
    for i, hd in enumerate(["Style", "Size", "Weight"]):
        txt(sl, ML + 0.1 + i * (tw / 3), y + 0.045, tw / 3 - 0.15, 0.2, hd, size=8,
            color=WHITE, bold=True, fit_check=False)
    for j, r in enumerate(rows):
        ry = y + 0.26 + j * 0.255
        rect(sl, ML, ry, tw, 0.255, fill=G50 if j % 2 == 0 else WHITE)
        rect(sl, ML, ry, tw, 0.007, fill=G200)
        for i, cell in enumerate(r):
            txt(sl, ML + 0.1 + i * (tw / 3), ry + 0.045, tw / 3 - 0.15, 0.2, cell,
                size=8.2, color=NAVY if i == 0 else G700, bold=(i == 0),
                fit_check=False)
    txt(sl, ML, y + 0.26 + 9 * 0.255 + 0.1, tw, 0.42,
        "Headings: line-height about 1.08, letter-spacing -0.02em. Body line-height 1.5. Hero and proof "
        "numerals are oversized, 72 to 120pt in print. Bricolage Grotesque and IBM Plex Sans are not "
        "brand fonts.",
        size=8.2, color=G600, line_spacing=1.3)
    # right column
    rx = ML + tw + 0.30
    rw = CW - tw - 0.30
    cards_x = [
        ("The eyebrow technique, a signature",
         "A small uppercase label above a section heading. 13px, weight 700, letter-spacing 0.14em, "
         "colour #143DFF, margin-bottom 18px. Use it on every major section and every content slide."),
        ("Spacing scale, these steps only",
         "2 · 4 · 8 · 12 · 16 · 20 · 24 · 32 · 40 · 48 · 60 · 80 · 96 · 120 px. Component rhythm favours "
         "8, 12, 16, 24. Section padding favours 60, 80, 96, 120. Container max-width 1200px."),
        ("Border radius",
         "0 · 4 · 5 · 15 · 20 · 30 · 40 · 9999 px. Buttons and inputs 4 to 5. Chips and small cards 15. "
         "Cards and panels 20. Large surfaces and the footer band 30 to 40. Pills and avatars 9999."),
        ("Buttons",
         "Two shapes, four fills. Satoshi 700 at 16px, padding about 14px 26px, hover lift, focus ring "
         "3px #B1BDFF at 2px offset. Rectangular for primary page actions, pill for content and "
         "navigation actions. On navy, the primary call to action is a white button with dark text."),
    ]
    cy = top + 0.28
    for t, b in cards_x:
        h = 0.94
        rect(sl, rx, cy, rw, h, fill=G50, radius=20)
        rect(sl, rx, cy, 0.055, h, fill=BLUE)
        txt(sl, rx + 0.25, cy + 0.15, rw - 0.45, 0.22, t, size=10.2, color=NAVY,
            bold=True)
        txt(sl, rx + 0.25, cy + 0.42, rw - 0.45, h - 0.56, b, size=8.2, color=G700,
            line_spacing=1.26)
        cy += h + 0.14
    # button samples
    b1 = rect(sl, rx, cy + 0.02, 1.55, 0.34, fill=BLUE, radius=5)
    txt(sl, rx, cy + 0.10, 1.55, 0.2, "Get in touch", size=9, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, fit_check=False)
    b2 = rect(sl, rx + 1.70, cy + 0.02, 1.75, 0.34, fill=BLACK, radius=5)
    txt(sl, rx + 1.70, cy + 0.10, 1.75, 0.2, "Request a Demo", size=9, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER, fit_check=False)
    b3 = rect(sl, rx + 3.60, cy + 0.02, 1.85, 0.34, fill=None, radius=9999,
              line=NAVY, line_w=1.2)
    txt(sl, rx + 3.60, cy + 0.10, 1.85, 0.2, "Read Articles  ›", size=9, color=NAVY,
        bold=True, align=PP_ALIGN.CENTER, fit_check=False)


page(s_type)


def s_dodont(sl):
    top = slide_head(sl, "03 · DESIGN SYSTEM", "Do, don't, and the print variant",
                     "Land the design system, do not reinvent it. Every deviation costs a rebuild later.")
    two_col(
        sl,
        ("Do", [
            "Keep electric blue as one sharp accent",
            "Lead proof zones with oversized numerals, because the number is the hero",
            "Use the navy radial glow with grain or a measurement-grid texture",
            "Prefer product imagery: a real 3D body-scan render, the guided-capture phone user interface, the Admin Panel in a browser frame",
            "Use the type scale and spacing steps exactly",
            "Use the eyebrow technique on every major section",
            "White call-to-action button on navy",
            "Honour reduced-motion preferences",
        ], BLUE),
        ("Don't", [
            "Spread #143DFF across large fills",
            "Use flat navy blocks with no depth",
            "Use an icon where a product asset fits",
            "Invent a size or a radius that is off-scale",
            "Put light text on imagery without a scrim",
            "Introduce a font other than Satoshi",
            "Resurrect #2962FF or Inter",
            "Rebuild the triangle logo mark by eye. Clearspace and minimum sizes are not in the Figma export, so request the brand-mark asset kit from Vadim",
        ], NAVY),
        top, height=2.72, lfill=BLUE_50, rfill=G50, bsize=8.8, tsize=12,
    )
    cards(sl, [
        ("Print and collateral variant",
         "For a two-pager or a printed leave-behind, the print art direction uses a marginally warmer "
         "navy #0A1338, ink #0B0B0C and muted #5D6070, with the same blue tints. On web, use the "
         "canonical values. Electric blue #143DFF is identical everywhere. Embed or outline Satoshi in "
         "print PDFs."),
        ("Footer, the link section",
         "A full-width navy radial-glow band with rounded top corners at about 30 to 40px. Centred "
         "heading in Display or H1 Satoshi Bold white: \"Let's talk\", \"Request a Demo\" or \"Unlock "
         "Body Data\". A white rectangular \"Get in touch\" button. Two centred outlined trust pills: "
         "HIPAA Compliant, GDPR."),
        ("Accessibility and motion",
         "Body text meets AA contrast minimum. Never place light text on light imagery without a scrim. "
         "Visible focus ring #B1BDFF at 3px with 2px offset. Scroll-reveal with stagger, transitions "
         "about 0.15s ease, no bounce or overshoot, and always honour prefers-reduced-motion."),
    ], top + 2.90, cols=3, height=1.72, tsize=10.2, bsize=8.2, fill=WHITE,
        line=G200)


page(s_dodont)


# ============================================================ 04 USE CASES
def s_div4(sl):
    divider(sl, "04", "Use-case library",
            "Eleven health segments, each with the same fields. The working reference, and the section Nika will open most often.")


page(s_div4, dark=True)


def s_uc_intro(sl):
    top = slide_head(sl, "04 · USE-CASE LIBRARY", "How to read this section, and the boundary rule",
                     "Sources: audience.md for the hook and the what-not-to-say, icp-detail.md for buyers and thresholds, proof-points.md for every number, case-studies/ for every customer name.")
    banner(sl, "RESPECT THE VERTICAL BOUNDARIES",
           "Telehealth copy does not bleed into online-pharmacy compliance. Fitness copy does not bleed "
           "into GLP-1 clinical workflows. Wellness copy does not bleed into clinical decisioning. Each "
           "vertical page and each deck holds its own lane, otherwise two assets compete for the same "
           "query and the same buyer, and neither reads as written by someone who knows the market.",
           top - 0.32, height=1.06, tsize=12)
    txt(sl, ML, top + 0.86, CW, 0.24,
        "PROOF-AVAILABILITY REALITY, READ THIS BEFORE PROMISING A DECK", size=8.6,
        color=BLUE, bold=True)
    table(sl, ["Tier", "Segments", "What that means for collateral"], [
        ["Named customer with a number",
         "Telehealth and GLP-1 (Yazen), online pharmacy (UK Meds)",
         "A case card and a vertical proof point are possible, subject to naming permission"],
        ["Use-case file and market sizing, no customer",
         "Insurance, health plans and wellness, bariatric, occupational health, clinical trials, digital fitness",
         "The vertical case slot stays empty and is recorded. Build on workflow, governance and the product proof set"],
        ["No use-case file, no proof",
         "Plastic surgery, breast cancer-related lymphedema (BCRL)",
         "Needs proof and needs a use-case file. Flag to Vadim before any commitment"],
    ], top + 1.14, widths=[0.24, 0.34, 0.42], row_h=0.56, bold_first=True)
    cards(sl, [
        ("Telehealth and GLP-1 are split here",
         "audience.md treats them as one segment with two lenses. This playbook splits them into 4.1 and "
         "4.2 because they need different decks. They share proof, and they must not contradict each other."),
        ("Market sizing is internal only",
         "The total and serviceable addressable market ranges in proof-points.md come from the Brainstorm "
         "deck and are marked illustrative. Use them for prioritization and internal business cases. Never "
         "in customer-facing collateral."),
        ("Each segment carries eight fields",
         "Ideal customer profile and buyer titles, core pain, the FitXpress fit, the hero outcome, the proof "
         "to use, what NOT to say, deck localization, and the landing-page position."),
    ], top + 3.26, cols=3, height=1.50, tsize=10.2, bsize=8.2, fill=G50,
        accent_bar=True)


page(s_uc_intro)


def s_uc_matrix(sl):
    top = slide_head(sl, "04 · USE-CASE LIBRARY", "The eleven segments at a glance",
                     "Priority follows proof. Where the proof column says \"needs proof\", the case card and the testimonial slot stay empty and are recorded.")
    table(sl, ["#", "Segment", "Revenue floor", "Proof status", "Collateral state today"], [
        ["4.1", "Telehealth and virtual care", "$2M+", "Yazen, 34,000 scans in 2025",
         "Benchmark landing page live. Extend rather than duplicate"],
        ["4.2", "Weight loss and GLP-1 programs", "$2M+", "Yazen, shares the 4.1 set",
         "Needs its own page, linked to and from telehealth"],
        ["4.3", "Online pharmacy and remote prescribing", "$2M+", "UK Meds, 7,500 scans in 2025",
         "/for-bmi-verification/ is the priority rewrite"],
        ["4.4", "Insurance underwriting", "$5M+", "Needs proof",
         "Deck written, and it is the structural exemplar. No page"],
        ["4.5", "Health plans and employer wellness", "$5M+", "Needs proof", "No page"],
        ["4.6", "Bariatric and metabolic clinics", "Mid to large", "Needs proof",
         "Use-case one-pager for business development first"],
        ["4.7", "Occupational health and screening", "Mid to large", "Needs proof",
         "Legal read on the framing before any public page"],
        ["4.8", "Clinical trials, CROs and pharma", "Global", "Needs proof",
         "Published article covers the top of the funnel"],
        ["4.9", "Connected and digital fitness", "$1M+", "Needs proof",
         "Page live, and it needs a FAQ, schema and a URL fix"],
        ["4.10", "Plastic surgery clinics, Turkey first", "$1M+", "No proof, no use-case file",
         "Blocked. No deck until the use-case file exists"],
        ["4.11", "BCRL and oncology survivorship", "$2M+", "No proof at all",
         "Blocked pending proof and a Whitney review"],
    ], top, widths=[0.045, 0.245, 0.11, 0.22, 0.38], row_h=0.335, bold_first=False)
    txt(sl, ML, top + 4.14, CW, 0.24,
        "Seven of eleven segments carry no citable customer. That is the single biggest constraint on "
        "what marketing can promise, and it is item 12 in section 9.",
        size=9, color=G600, italic=True)


page(s_uc_matrix)


# --- per-segment slides
def segment_slide(num, name, hero, hero_note, icp, buyers, pain, fit, proof, never,
                  deck, page_note, hero_blocked=False):
    def fn(sl):
        eyebrow(sl, f"04 · USE-CASE LIBRARY  ·  {num}")
        txt(sl, ML, 0.70, CW, 0.44, name, size=26, color=NAVY, bold=True,
            line_spacing=1.06)
        # hero band
        if hero_blocked:
            rect(sl, ML, 1.20, CW, 0.66, fill=G100, radius=20, line=G300, line_w=1.2)
            txt(sl, ML + 0.24, 1.31, 2.3, 0.18, "HERO OUTCOME", size=8, color=G500,
                bold=True, fit_check=False)
            txt(sl, ML + 0.24, 1.50, CW - 0.5, 0.30, hero, size=11.5, color=G700,
                bold=True, line_spacing=1.14)
        else:
            banner(sl, "HERO OUTCOME", hero, 1.20, height=0.66, tsize=13)
        txt(sl, ML, 1.92, CW, 0.20, hero_note, size=8.2, color=G600, italic=True)
        cards(sl, [
            ("Ideal customer profile", icp),
            ("Core pain", pain),
            ("The FitXpress fit", fit),
            ("Buyer titles", buyers),
            ("Deck localization", deck),
            ("Landing page", page_note),
        ], 2.20, cols=3, height=1.30, gap=0.18, tsize=9.6, bsize=7.9, fill=G50,
            accent_bar=True, pad=0.17)
        two_col(
            sl,
            ("Proof to use", proof, BLUE),
            ("What NOT to say", never, NAVY),
            5.00, height=1.96, lfill=BLUE_50, rfill=G100, bsize=8.2, tsize=10.5,
        )
    return fn


page(segment_slide(
    "4.1", "Telehealth and virtual care",
    "Verify body progress remotely to boost retention, reduce drop-off, and prove program ROI",
    "Shorter variant: Make body progress visible before members drop off.",
    "Virtual-first clinics, telehealth platforms, longitudinal care and remote patient monitoring "
    "programs, cardiometabolic platforms. Series A through public, 50 to 5,000 employees, $2M+ annual "
    "revenue. Typically 500+ active members with repeat check-ins. USA, Canada, UK, Germany, UAE, Australia.",
    "Founder or CEO · Chief Medical Officer or Medical Director · Head of Clinical Operations · Head of "
    "Member Engagement or Retention · Head of Product · CTO · Head of Outcomes and Program Insights.",
    "Remote programs run on repeat check-ins, and the data underneath is shaky. Self report and "
    "occasional progress photos are inconsistent, easy to skip and hard to compare across scans. "
    "Clinical teams cannot separate real change from measurement noise. Manual intake slows scaling.",
    "A structured intake and documentation layer inside the existing app flow. Two photos return 80+ "
    "measurements and body composition, and because variance stays < 1 cm, the trend a care team reads "
    "is change rather than noise.",
    ["Yazen, 34,000 scans in 2025, weight-loss management support. Naming subject to permission",
     "80+ measurements, under 45 seconds from two photos",
     "96 to 97% accuracy against expert manual measurement, typical absolute error 1.5 to 2.0 cm",
     "Repeatability < 1 cm, 95%+ consistency. Weight estimation plus or minus 3.5%",
     "HIPAA maintained, GDPR principles, encryption at rest and in transit, photos deleted immediately or within 30 days"],
    ["No diagnostic claims",
     "Do not position as a DEXA or calibrated-scale replacement",
     "Do not imply eligibility decisioning",
     "Keep separate from UK online-pharmacy Body Mass Index (BMI) compliance, unless the piece is explicitly the bridge"],
    "Problem slide on trust in the trend rather than one-time accuracy. Journey slide on the 30, 60 and "
    "90 day check-in cadence. Outputs slide member-facing. Accuracy slide leads with repeatability.",
    "/structured-body-data-for-telehealth-digital-health-programs/ already covers this vertical well. "
    "Extend rather than duplicate. Bottom of funnel: \"Book a demo\", with \"See sample outputs\" secondary.",
))

page(segment_slide(
    "4.2", "Weight loss and GLP-1 programs",
    "Make body progress visible between check-ins, helping programs keep members past day 90",
    "The source line used \"so\" as a benefit connector, which Asselya's rules ban. The rewrite above is the approved form.",
    "GLP-1 prescription platforms, metabolic and obesity platforms, coaching apps treating weight "
    "management, employer-sponsored metabolic health programs. $2M+ annual revenue. USA, UK, Germany, "
    "Australia, Canada, UAE.",
    "Founder or CEO · Chief Medical Officer · Head of Clinical Operations · Head of Member Engagement · "
    "Head of Outcomes · Growth and Revenue Operations.",
    "Members drop off when progress is not visible, and weight alone hides what matters because it does "
    "not separate fat from lean mass. GLP-1 prescribing needs reliable baseline and follow-up body "
    "metrics. Rising acquisition cost with weak engagement and early churn. Payers increasingly expect "
    "defensible outcomes.",
    "Repeat scans that make composition change visible between clinical assessment points, feeding both "
    "member engagement and outcomes reporting.",
    ["The same set as 4.1. Yazen is the anchor",
     "Lean-mass preservation tracking is a buyer KPI worth naming, framed as what the data supports rather than what the product proves",
     "Body composition outputs: Body Mass Index (BMI), basal metabolic rate (BMR), body fat percentage, lean mass, fat mass",
     "The target-weight 3D visualization as an engagement mechanic"],
    ["No claim about the clinical outcomes of GLP-1 therapy",
     "No eligibility or dosing implication",
     "No lean-mass-preservation efficacy claim. FitXpress documents the measurement, and the clinical interpretation belongs to the care team"],
    "Problem slide on day-90 drop-off. A dedicated slide on the progress and target-weight 3D "
    "visualization as an engagement mechanic. Outcomes slide on retention, adherence and outcomes "
    "reporting to payers.",
    "Its own page, linked from and to the telehealth page. Middle of funnel: \"Review the buyer "
    "checklist\". Bottom of funnel: \"Book a demo\".",
))

page(segment_slide(
    "4.3", "Online pharmacy and remote prescribing",
    "Verify BMI inside the order flow to cut fraud, speed approvals, and stay audit-ready",
    "The United Kingdom is the priority market for this segment.",
    "Online pharmacies and remote-prescribing platforms dispensing GLP-1 and weight-loss medication, "
    "needing to verify a Body Mass Index (BMI) threshold, typically 27 or above with a comorbidity or 30 "
    "or above, without a visit. 50 to 1,000 employees, $2M+ annual revenue. UK priority, plus "
    "international telehealth providers operating in the UK.",
    "Head of Compliance and Risk · Chief Medical Officer · Clinical Operations Director · Head of "
    "Clinical Governance · Founder or CEO · Product Manager · CTO.",
    "Patients misreport weight and BMI to qualify. Manual photo review is subjective and inconsistent. "
    "UK regulators are tightening scrutiny of remote GLP-1 eligibility. High intake volumes overload the "
    "clinical team, and incorrect verification carries prescribing, legal and reputational exposure.",
    "Verification inside the order flow, with a structured record behind it. Smart Scales compares "
    "self-reported weight against the estimate and flags a mismatch for clinician review. Integration "
    "Pattern B fits: photos in, eligibility validated server-side, body metrics never shown to the user, "
    "audit trail retained.",
    ["UK Meds, 7,500 scans in 2025, BMI verification for online pharmacy ordering. Naming subject to permission",
     "Weight estimation plus or minus 3.5% average error",
     "Under 45 seconds, two photos. Repeatability < 1 cm",
     "The audit and compliance posture from compliance.md, and Pattern B from tech-spec.md"],
    ["Never imply automated eligibility or prescribing decisions",
     "Never guarantee compliance. \"Supports compliant workflows\", never \"makes you compliant\"",
     "Never \"detects fraud\". Always \"flags discrepancies for review\"",
     "Name UK regulators only when the piece is genuinely about the UK market, and expand them at first use: Medicines and Healthcare products Regulatory Agency (MHRA), Care Quality Commission (CQC)"],
    "Problem slide on manual photo review and the regulatory tightening. A workflow slide showing "
    "Pattern B explicitly, including who reviews a flagged scan. Compliance and audit trail elevated to "
    "a deciding slide rather than a footnote.",
    "/for-bmi-verification/ exists at about 659 words with no FAQ and no cases. It is the priority "
    "rewrite, and the week-two deliverable. Bottom of funnel: \"Book a demo\".",
))

page(segment_slide(
    "4.4", "Insurance underwriting, life and disability",
    "Verify body metrics remotely to issue faster, cut rework, and strengthen auditability",
    "The shipped deck headline is \"Stronger evidence, faster underwriting\". This is the structural exemplar deck.",
    "Life and disability insurers, group benefits carriers, reinsurers, insurtech and digital "
    "distribution platforms, accelerated and digital underwriting teams. Large enterprise, typically "
    "5,000+ employees, $5M+ annual revenue. US core, UK and EU, Canada, UAE, Australia.",
    "Chief Underwriting Officer or VP Underwriting · Chief Medical Officer · Chief Risk Officer or Head "
    "of Risk and Analytics · Head of Digital Innovation · Population Health Director · CTO · Compliance "
    "and Fraud Prevention.",
    "Accelerated underwriting removed friction and created an evidentiary gap. Self-reported height, "
    "weight, build and Body Mass Index (BMI) can be incomplete, inconsistent or intentionally "
    "misrepresented, which affects risk classification, pricing accuracy and fraud exposure. The "
    "alternatives reintroduce the delay accelerated programs were designed to remove.",
    "A structured digital evidence layer: remote guided capture producing repeatable, auditable body "
    "data that supports underwriter review, applicant triage and fraud-prevention workflows, without "
    "becoming a decisioning engine.",
    ["No FitXpress insurance customer exists. NEEDS PROOF, flag to Vadim for a case card, a quote or a reference",
     "Build on 96 to 97% accuracy, error margin 1.5 to 2.0 cm, repeatability < 1 cm and 95%+ consistency, weight estimation plus or minus 3.5%, 80+ measurements, under 45 seconds",
     "The full governance set: encryption, configurable retention, no personal identifiers, logging of scan status, timestamps, quality flags and failure reasons",
     "External context may be cited with a named source and a link (CDC, Munich Re), and only after verification at source and external-use approval"],
    ["Never automated underwriting, never automatic fraud detection, never standalone decisioning",
     "Never employment screening",
     "Never \"best-in-class repeatability\", which the shipped deck currently contains and which is a banned word",
     "HIPAA applicability to life insurers varies, so do not assert it as given. Confirm per prospect"],
    "This is the exemplar deck, already written. Reuse its shape: alternatives and their limitations on "
    "the problem slide, five outcome cards with \"driven by\" and \"track with\", a five-stage "
    "underwriting journey, and the integration boundary slide that keeps the carrier in control.",
    "Does not exist. Build after the case-evidence gap is resolved. Middle of funnel: \"Review the "
    "evidence-workflow checklist\". Bottom: \"Talk to 3DLOOK about your underwriting workflow\".",
))

page(segment_slide(
    "4.5", "Health plans and employer wellness",
    "Verify wellness progress remotely to reduce disputes, boost participation, and improve program reporting",
    "Rewards verification. Fairness and consistency is what this buyer actually worries about.",
    "Health insurers running wellness incentive programs, Medicare Advantage providers, self-insured "
    "employers, employer benefit and wellness platforms, population-health vendors. Large enterprise, "
    "$5M+ annual revenue. USA, Canada, UK, Germany, UAE, Australia.",
    "CHRO · Head of Wellness or Director of Wellbeing · VP Population Health · Chief Medical Officer · "
    "Health Plan Operations Director · Compliance and Risk · Wellness Program Administrator.",
    "Manual wellness verification is an administrative load and self report is unreliable. Hybrid and "
    "remote workforces complicate onsite screening. Complex verification cuts participation, delayed "
    "verification weakens the reward's motivational effect, and disputes and inconsistent submissions "
    "undermine fairness and trust.",
    "Remote, standardized, audit-ready verification of a biometric milestone, with fraud-prevention "
    "support and consistent treatment across a distributed population.",
    ["No customer. NEEDS PROOF, flag to Vadim",
     "Bupa Rewards is noted internally as a public reference for what a good program looks like. It is not a 3DLOOK customer and must never be presented as one",
     "Build on the product proof set and the governance set",
     "Participation-rate and dispute-rate figures do not exist. Use them as \"track with\" metric categories, never as claimed results"],
    ["No medical or diagnostic claims",
     "Never imply the reward or the eligibility decision is automated",
     "Keep every benefit claim soft: supports, may reduce, can help"],
    "Outcomes slide on participation, dispute volume and cost per validated check-in. A fairness and "
    "consistency slide. Compliance and data-governance slide elevated.",
    "Does not exist. Middle of funnel: \"See how remote verification fits an incentive program\". "
    "Bottom: \"Book a demo\".",
))

page(segment_slide(
    "4.6", "Bariatric and metabolic clinics",
    "Pre-qualify patients remotely to reduce wasted consults, speed pre-auth, and improve conversion to procedures",
    "Pre-qualification. This vertical is better served first by a use-case one-pager for business development.",
    "Bariatric surgery centres, hospital programs, multi-site surgical networks, metabolic and obesity "
    "clinics. Mid to large hospital systems and surgical-centre chains. US and EU.",
    "Director of Operations · Medical Director · VP Patient Access · COO for multi-site networks.",
    "Long consult waitlists, a high rate of late-stage disqualification that wastes consult slots, "
    "pre-authorization backlogs, and staff time consumed per pre-authorization.",
    "Remote pre-qualification before the consult, producing a structured record that supports "
    "pre-authorization documentation.",
    ["No customer. NEEDS PROOF, flag to Vadim. Product proof set only",
     "KPI categories to offer as \"track with\": consult-to-procedure conversion, pre-authorization cycle time, wasted consult rate, staff time per pre-authorization",
     "use-cases/fx-bariatric-pre-auth.md exists as the starting point"],
    ["Never a surgical candidacy determination",
     "Never a claim about anaesthetic or surgical risk",
     "Never guaranteed payer approval"],
    "Journey slide from enquiry through remote pre-screen to consult booking. Operational slide on "
    "throughput and staff time.",
    "Later. Start with the use-case one-pager for business development, not a page.",
))

page(segment_slide(
    "4.7", "Occupational health and pre-employment screening",
    "Standardize screening intake remotely to increase throughput, reduce rescreens, and speed clearance decisions",
    "This is the vertical where a claims slip is most expensive. Get a legal read on the framing before any public page.",
    "Occupational health providers, workforce screening vendors, workers' compensation and absence "
    "administrators, large multi-site employers. Mid to large. US and EU.",
    "VP Operations or COO · Chief Medical Officer · Director of Clinical Services · Head of Occupational "
    "Health.",
    "Throughput per clinic per day, rescreen and rework rates, time to clearance, and audit readiness "
    "across sites that each do it slightly differently.",
    "Standardized remote intake that reduces variability between sites and reviewers, with a documented "
    "record.",
    ["No customer. NEEDS PROOF, flag to Vadim. Product proof set only",
     "Repeatability and the multi-site standardization argument carry this vertical",
     "The governance set: audit trails, timestamps, quality flags, failure reasons"],
    ["Never a fitness-for-duty determination, never clearance decisioning, never a hiring decision, never a medical assessment",
     "The framing is intake standardization and documentation. The determination is the provider's",
     "Employment screening is explicitly excluded from what FitXpress claims"],
    "Multi-site standardization slide. Audit-readiness slide. Journey slide from candidate invitation to "
    "clearance review.",
    "Later. Get a legal read on the framing before any public page in this vertical.",
))

page(segment_slide(
    "4.8", "Clinical trials, CROs and pharma sponsors",
    "Standardize anthropometrics across sites and reduce visit burden to improve data quality and retention",
    "Contract Research Organization (CRO). The published clinical-trials article is the best existing model for compliance scoping.",
    "Contract Research Organizations (CROs), pharma sponsors running metabolic and obesity trials, "
    "academic research networks, decentralized clinical trial (DCT) platforms. Global.",
    "Director of Clinical Operations · VP Decentralized or Hybrid Trials · Head of Site Management · "
    "Director of Data Management.",
    "Site-to-site variability in anthropometric measurement, coordinator time per visit, retention and "
    "dropout, screen-failure rate, and audit findings on measurement variability across sites.",
    "Standardized anthropometric capture with a timestamped, structured record that reduces visit burden "
    "and inter-site variability.",
    ["No customer. NEEDS PROOF, flag to Vadim. Product proof set, and the repeatability figures specifically, which are the argument here",
     "The published article clinical-trials-anthropometric-measurement.md is the best existing model. Read its known_issues frontmatter to avoid repeating its M1 and M2 slips",
     "Repeatability by measurement: chest 0.60 cm, waist 0.89 cm, knee 0.12 cm, ankle 0.07 cm"],
    ["Never endpoint validation",
     "Never a claim of protocol or regulatory compliance",
     "Never a Good Clinical Practice (GCP) or Food and Drug Administration (FDA) qualification claim",
     "Expand every regulator and standard at first use, which is exactly where this vertical's drafts slip"],
    "Variability-across-sites problem slide. Data-management and export slide, since Electronic Data "
    "Capture (EDC) and electronic Clinical Outcome Assessment (eCOA) integration questions will come up. "
    "Audit-trail slide.",
    "Later. The published article already covers the top of the funnel for this vertical.",
))

page(segment_slide(
    "4.9", "Connected and digital fitness",
    "Give members visible body progress, and give the product a reason to be renewed",
    "The tone here is lighter and less clinical than every other segment in this library.",
    "Connected fitness platforms, digital fitness and coaching apps, virtual personal training, fitness "
    "subscription platforms, corporate wellness and fitness platforms. Series B to public, $1M+ annual "
    "revenue, strong recurring-subscription focus. US, UK, Canada, Germany, UAE, Australia, Nordics.",
    "Founder or CEO · Chief Product Officer · Head of Growth · VP User Engagement or Retention · CTO · "
    "Product Manager.",
    "Users lose motivation without visible progress, churn after onboarding is high, personalization is "
    "limited to surveys and stated goals, and the market is crowded and competing on user experience. "
    "Rising acquisition cost demands stronger engagement.",
    "Visible transformation and body-data personalization as an engagement and premium-tier mechanic.",
    ["proof-points.md holds no named fitness customer. NEEDS PROOF, flag to Vadim",
     "Names circulating in the sales deck and the telehealth page draft are not in proof-points.md or case-studies/, so they are unusable until Vadim adds them with permission",
     "Build on 80+ measurements, body composition outputs, the target-weight 3D visualization, side-by-side longitudinal comparison, under 45 seconds, repeatability < 1 cm"],
    ["No medical, diagnostic or clearance language at all",
     "Do not blur into GLP-1 clinical workflows or wellness-rewards verification",
     "This is the one vertical where the clinical register actively hurts the pitch"],
    "Engagement and retention outcomes slide. A visual-progress slide carrying more weight than the "
    "accuracy slide. Shorter deck overall, 12 to 15 slides.",
    "/fitxpress/for-connected-and-digital-fitness/ exists at about 1,352 words with no FAQ and no "
    "schema. Two jobs: add the FAQ block with FAQPage schema, and resolve the URL, which sits on a "
    "/fitxpress/ path level that redirects. Coordinate with Vadim.",
))

page(segment_slide(
    "4.10", "Plastic surgery clinics, Turkey first",
    "Pre-qualify international patients before they book a flight",
    "Turkey is the priority geography for medical tourism, then USA, Canada, UK, Germany, UAE, Australia.",
    "Plastic surgery clinics and aesthetic medicine centres. Large, 10+ surgeons and 50+ procedures a "
    "month, or mid-sized, 3 to 10 surgeons and 20 to 50 procedures a month. $1M+ annual revenue, "
    "premium pricing.",
    "Clinic Owner or Director · Plastic Surgeon · Project Manager or Technical Lead.",
    "High drop-off between consultation and surgery because patients are disqualified at the last moment "
    "on Body Mass Index (BMI). Wasted in-person consult slots. No reliable remote pre-screen, and "
    "ordinary scales are not enough. Fat distribution and composition affect surgical outcome but are "
    "poorly measured before the visit.",
    "Remote pre-qualification on BMI and body composition before travel or a consult, and better "
    "structured input for surgical planning.",
    ["No customer, no use-case file, and no vertical proof",
     "NEEDS PROOF and NEEDS A USE-CASE FILE. use-cases/fx-plastic-surgery.md does not exist. Flag to Vadim",
     "Do not build a deck until at least the use-case file exists, because the vertical-context slot cannot be filled from any current source: who signs off, how the medical-tourism funnel runs, what the coordinator does with a flagged scan"],
    ["Never a diagnostic claim",
     "Never an anaesthesia-risk or surgical-risk claim",
     "Never a replacement for surgeon assessment",
     "The framing is support for pre-screening and planning"],
    "Not yet. When it is unblocked: a medical-tourism funnel journey slide, and a coordinator-workflow "
    "slide, since the coordinator is the daily user.",
    "Not yet.",
))

page(segment_slide(
    "4.11", "BCRL detection and monitoring, oncology survivorship",
    "Not yet writable. Any outcome line here would rest on a volumetric-asymmetry tracking claim that no current proof point supports",
    "Breast cancer-related lymphedema (BCRL). Blocked pending proof and a Whitney review.",
    "Oncology survivorship platforms, breast-cancer aftercare clinics, remote patient monitoring "
    "providers, cancer care networks, lymphedema and rehabilitation teams, cancer treatment centres and "
    "hospital systems. $2M+ annual revenue. USA, Canada, Germany, UAE, Australia.",
    "Chief Medical Officer · Oncology Program Director · Survivorship Care Leader · Remote Patient "
    "Monitoring Director · Rehabilitation and Lymphedema Specialist · Digital Health Product Manager · CTO.",
    "BCRL is underdiagnosed early. Tape measurement is inconsistent and hard to reproduce at home. "
    "In-clinic monitoring does not scale and long-term self-monitoring adherence is weak. Subtle "
    "volumetric and asymmetry changes are hard to capture, and documentation is slow and manual.",
    "Remote longitudinal monitoring support and reproducible digital body records, with engagement "
    "through visual 3D comparison.",
    ["NONE. This segment has no proof point and no use-case file",
     "icp-detail.md records explicitly that the clinically relevant metric, volumetric asymmetry tracking accuracy, is absent from proof-points.md",
     "Flag to Vadim. No deck, no page, no one-pager until a validated metric exists and Whitney has reviewed the framing"],
    ["Never detect or diagnose lymphedema or BCRL",
     "Never replace clinical assessment",
     "\"Supports monitoring workflows\" is the only safe frame",
     "This is the most compliance-sensitive vertical in the whole portfolio"],
    "Blocked. Recommendation: treat this as a research request to product before it is treated as a "
    "marketing request.",
    "Blocked pending proof and a Whitney review.",
    hero_blocked=True,
))


# ============================================================ 05 DECKS
def s_div5(sl):
    divider(sl, "05", "Deck build playbook",
            "The first of the two production lines. A fixed spine of product truth, plus localization.")


page(s_div5, dark=True)


def s_deck_principle(sl):
    top = slide_head(sl, "05 · DECK BUILD PLAYBOOK", "The principle, and the lengths",
                     "A vertical deck is not written from scratch. Building it any other way is how two decks end up quoting two different numbers for the same thing.")
    # 60/40 bar
    bw = CW
    rect(sl, ML, top, bw * 0.6, 0.62, fill=NAVY, radius=None)
    rect(sl, ML + bw * 0.6, top, bw * 0.4, 0.62, fill=BLUE, radius=None)
    txt(sl, ML + 0.22, top + 0.12, bw * 0.6 - 0.44, 0.42,
        "About 60%: a fixed spine of product truth, identical in every vertical",
        size=12, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, fit_check=False)
    txt(sl, ML + bw * 0.6 + 0.22, top + 0.12, bw * 0.4 - 0.44, 0.42,
        "About 40%: localization", size=12, color=WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE, fit_check=False)
    txt(sl, ML, top + 0.78, CW, 0.46,
        "The exemplar is fitxpress-insurance-underwriting-deck-copy.md. It is 29 slides built by "
        "localizing the telehealth deck, and it documents its own text-fit review and its own open items "
        "at the end. Copy that discipline, and copy its structure. Do not copy its two banned-word slips, "
        "which are items 9 and 10 in section 9.",
        size=10, color=G700, line_spacing=1.3)
    txt(sl, ML, top + 1.36, CW, 0.24, "DECK LENGTHS", size=8.6, color=BLUE, bold=True)
    table(sl, ["Deck", "Slides", "Use"], [
        ["Demo deck", "19 core plus appendix, 25 to 29 total", "The main vertical asset, sent after a demo"],
        ["First-call deck", "10 to 12", "Slides 1 to 8, plus compliance and the close"],
        ["Leave-behind or two-pager", "2 pages", "Hero, outcomes, how it works, proof, compliance, contact"],
        ["Conference or keynote", "8 to 10", "Problem, why now, workflow, proof, close"],
    ], top + 1.64, widths=[0.24, 0.26, 0.50], row_h=0.36, bold_first=True)
    banner(sl, None,
           "About 60% of a vertical deck is product truth that must be identical everywhere. Retyping a "
           "fixed slide is how a number drifts, and a drifted number is the first thing a diligence reader "
           "catches.",
           top + 3.42, height=0.72, tsize=12.5)


page(s_deck_principle)


def s_deck_rules(sl):
    top = slide_head(sl, "05 · DECK BUILD PLAYBOOK", "Deck design rules, and roadmap discipline",
                     "Land the design system, do not reinvent it. And promise only what has shipped.")
    cards(sl, [
        ("Deck design rules",
         "Navy #050F40 with the radial glow on the cover, section dividers, proof bands and the close. "
         "White on content slides. Electric blue #143DFF on one thing per slide, usually the number. "
         "Satoshi throughout, with the eyebrow label on every content slide. Proof slides lead with an "
         "oversized numeral. Prefer the three carrier assets over icon grids. Cards at 20px radius, chips "
         "at 15px. Confidentiality footer on every slide with the current month and year."),
        ("One text rule that saves a rebuild",
         "Write to the layout you have. If the replacement headline is much longer than what it replaces, "
         "shorten the headline rather than shrinking the type. Roughly 15% is the tolerance in the "
         "text-fit review before a layout change is needed."),
        ("Roadmap discipline",
         "Only what is shipped or in progress. Safe today: the clothing detector (available in the "
         "Application Programming Interface), anti-fraud live capture (in progress, not released), audit "
         "identifiers and logs (in release). Do not promise dates for height detection, Smart Scale "
         "accuracy improvements, the landmark-detector model upgrade, or full Software Development Kit "
         "branding. Confirm with Max Kucherenko before any roadmap slide ships."),
    ], top, cols=3, height=2.36, tsize=11, bsize=8.6, accent_bar=True)
    banner(sl, "WHEN IN DOUBT ON A ROADMAP QUESTION",
           "\"I don't want to give you a date I can't stand behind. Let me check with our product team and "
           "get back to you today.\" That answer costs nothing. A promised date that slips costs the deal.",
           top + 2.56, height=0.90, tsize=12.5)


page(s_deck_rules)


def _core_rows(a, b):
    rows = [
        ["1", "Cover", "L", "Outcome headline in the vertical's language, three proof points (2 smartphone photos, results in under a minute, no specialized hardware), the \"Built for [vertical]\" label. Highly visual, minimal text"],
        ["2", "The industry problem", "L", "Why existing approaches fail in this vertical. Four alternative-method blocks, each with what it is good at and its limitation. Close with a \"The result\" statement naming the business consequence"],
        ["3", "Business outcomes", "L", "Five outcome cards. Each: name, one sentence, two \"Driven by\" bullets, two or three \"Track with\" metrics. \"Track with\" names metric categories, never claimed results. Testimonial slot stays empty unless an approved quote exists"],
        ["4", "How it works", "F", "Four stages: user inputs, guided photo capture with pose and clothing validation, processing with photo deletion, structured return to your platform"],
        ["5", "What each scan returns", "F", "80+ measurements, body composition, 3D model, progress comparison, validation outputs (pose quality, clothing classification, Smart Scales mismatch flag)"],
        ["6", "Accuracy and repeatability, scoped", "F", "The reframe to \"accurate enough for which decision?\", the four conditions, then the figures from proof-points.md. Lead with repeatability for any longitudinal vertical. Never \"independent\", \"validated\", \"third-party\" or \"best-in-class\""],
        ["7", "Comparison", "L", "Compare by role. FitXpress against the alternative this buyer actually uses: self report, in-clinic DEXA or BIA, consumer photo apps, manual review. Never a clean sweep across every row"],
        ["8", "Where it fits in the workflow", "L", "The vertical's own five-stage journey, with the FitXpress step marked and the human decision point marked"],
        ["9", "The core workflow feature", "L", "The one capability this vertical buys: progress comparison, disclosure verification, eligibility pre-check, or inter-site standardization"],
        ["10", "The end-user experience", "L", "Adoption and completion in this vertical's flow. Where outputs are internal only (pharmacy Pattern B, underwriting), say so plainly and keep the slide short"],
        ["11", "Built for [vertical] operations", "L", "The operational requirements this buyer will ask about: volumes, review queues, exception handling, reporting"],
        ["12", "Reliability in production", "F", "Uptime posture, consistency across body types, lighting and capture environments, error handling"],
        ["13", "Built-in quality checks", "F", "Real-time pose validation, clothing detection, live capture, Smart Scales mismatch flag. State the limit of each control"],
        ["14", "Integration architecture", "F", "The \"we provide / you build\" boundary, the two integration patterns, the interface and the three kits. The boundary slide is what makes the buyer feel in control"],
        ["15", "Admin Panel", "F", "Centralized view and management of scan results, for teams that do not want to build a dashboard"],
        ["16", "What the team sees", "L", "Member-facing, internal-only, or both. Real white-label screenshots where permission allows"],
        ["17", "Security, privacy, compliance", "F+L", "The data-flow diagram, the controls list, HIPAA and GDPR posture. Localize only the opening sentence to the vertical's regulatory context"],
        ["18", "Pricing", "F", "The public price signal and \"no integration fee\". Never the internal per-request table. Check /pricing/ on the day the deck ships"],
        ["19", "How teams evaluate this", "F", "Three stages: demo evaluation, decision discussion, contracting. Ends on a concrete next step"],
    ]
    return rows[a:b]


def s_core_a(sl):
    top = slide_head(sl, "05 · DECK BUILD PLAYBOOK", "The 19-slide core, slides 1 to 10",
                     "F is fixed product truth and changes only when proof-points.md changes. L is localized per vertical. Every F-slide edit needs Vadim.")
    table(sl, ["#", "Slide", "F/L", "What goes on it"], _core_rows(0, 10),
          top, widths=[0.035, 0.20, 0.045, 0.72], row_h=0.44)


page(s_core_a)


def s_core_b(sl):
    top = slide_head(sl, "05 · DECK BUILD PLAYBOOK", "The 19-slide core, slides 11 to 19",
                     "Appendix, slides 20 onward: integration examples, the customer-controlled wrapper, capture, post-scan processing, output interface, internal-only consumption, About 3DLOOK, contacts.")
    y = table(sl, ["#", "Slide", "F/L", "What goes on it"], _core_rows(10, 19),
              top, widths=[0.035, 0.20, 0.045, 0.72], row_h=0.42)
    banner(sl, None,
           "Copy the fixed slides verbatim from the last approved deck. Do not retype them. Retyping is "
           "how a number drifts, and a drifted number is the first thing a diligence reader catches.",
           y + 0.24, height=0.66, tsize=12)


page(s_core_b)


def s_localize(sl):
    top = slide_head(sl, "05 · DECK BUILD PLAYBOOK", "How to localize, in order",
                     "The order matters. Steps 6 and 7 are separate passes after writing, because policing your own tells while drafting does neither job well.")
    items = [
        ("Fix the segment", "Open section 4. Read the segment's hook and its \"what NOT to say\" before writing a line."),
        ("Check the proof tier", "Is there a named customer with a number? If not, the case card and the testimonial slot stay empty and get recorded. Never fill them from an adjacent vertical."),
        ("Write slide 2 from the buyer's alternatives", "Not from our features. What does this team do today, and what specifically breaks?"),
        ("Write slide 3 from the buyer's KPI list", "From icp-detail.md. Outcomes are hedged: supports, may reduce, can help. \"Track with\" is a measurement category."),
        ("Rewrite only the L slides", "Copy the F slides verbatim from the most recent approved deck. Every F-slide edit is a change to product truth and needs Vadim."),
        ("Run the guardrails pass", "The 11 principles plus rules M1 and M2. As its own step, not while drafting."),
        ("Run the AI-tells sweep", "detect-ai-tells.py with --channel page. Fix every hard fail and house-rule violation, then read it yourself and answer what the detector cannot."),
        ("Write the text-fit review", "Slide by slide, is the new copy longer than what it replaces? The exemplar deck's review is the model, and roughly 15% is the tolerance."),
        ("Write the open items block", "Every unresolved number, every unapproved quote, every claim needing Whitney. Per guardrail 11, flag rather than decide."),
        ("Route it", "Vadim for numbers, Whitney for regulated claims, Asselya for editorial, Katerina for competitive framing."),
    ]
    cols, gap = 5, 0.18
    w = (CW - gap * (cols - 1)) / cols
    h = 2.34
    for i, (t, b) in enumerate(items):
        r, c = divmod(i, cols)
        x = ML + c * (w + gap)
        y = top + r * (h + 0.20)
        rect(sl, x, y, w, h, fill=G50, radius=20)
        rect(sl, x, y, w, 0.05, fill=BLUE, radius=None)
        cir = rect(sl, x + 0.18, y + 0.20, 0.30, 0.30, fill=BLUE, radius=9999)
        try:
            cir.adjustments[0] = 0.5
        except Exception:
            pass
        txt(sl, x + 0.18, y + 0.255, 0.30, 0.22, str(i + 1), size=10, color=WHITE,
            bold=True, align=PP_ALIGN.CENTER, fit_check=False)
        txt(sl, x + 0.18, y + 0.58, w - 0.36, 0.46, t, size=10, color=NAVY, bold=True,
            line_spacing=1.14)
        txt(sl, x + 0.18, y + 1.06, w - 0.36, h - 1.22, b, size=7.8, color=G700,
            line_spacing=1.24)


page(s_localize)


# ============================================================ 06 PAGES
def s_div6(sl):
    divider(sl, "06", "Landing-page build playbook",
            "The second production line. Four gates, seventeen slots, and the FAQ block that earns citations.")


page(s_div6, dark=True)


def s_gates(sl):
    top = slide_head(sl, "06 · LANDING-PAGE PLAYBOOK", "Route it through the pipeline",
                     "Pages on 3dlook.ai go through the page-builder skill. Nika owns the brief, the copy decisions and the final quality call. The pipeline enforces the parts that are easy to skip under deadline.")
    gates = [
        ("G-I", "Blocks the page existing at all",
         "A use-case file, two or more publishable cases from that vertical, real demand, five facts absent from the parent page, and the 60% uniqueness rule"),
        ("G-A", "Blocks writing",
         "Placement, URL, cannibalization against existing pages, and the Search Console baseline all settled"),
        ("G-T", "Blocks publishing, technically",
         "Schema, canonical, Yoast title and description, breadcrumbs, a single H1"),
        ("G-J", "Blocks publishing, on quality",
         "A blind judge in a fresh subagent, 100-point page scorecard, threshold 85, maximum three rounds. Publishing below 85 without flagging it is not allowed"),
    ]
    n, arrow = 4, 0.28
    w = (CW - arrow * 3) / 4
    h = 1.66
    for i, (g, blocks, test) in enumerate(gates):
        x = ML + i * (w + arrow)
        rect(sl, x, top, w, h, fill=G50, radius=20)
        rect(sl, x, top, w, 0.05, fill=BLUE, radius=None)
        txt(sl, x + 0.20, top + 0.20, w - 0.4, 0.32, g, size=19, color=BLUE, bold=True,
            fit_check=False)
        txt(sl, x + 0.20, top + 0.56, w - 0.4, 0.28, blocks, size=9.4, color=NAVY,
            bold=True, line_spacing=1.16)
        txt(sl, x + 0.20, top + 0.92, w - 0.4, h - 1.10, test, size=8.0, color=G700,
            line_spacing=1.26)
        if i < 3:
            txt(sl, x + w, top + h / 2 - 0.14, arrow, 0.28, "›", size=17,
                color=BLUE_300, bold=True, align=PP_ALIGN.CENTER, fit_check=False)
    banner(sl, "READ THIS BEFORE PROMISING A VERTICAL PAGE",
           "G-I currently blocks almost every FitXpress vertical, because it wants two publishable cases "
           "and every FitXpress vertical has at most one. The options are a second case, an approved "
           "reference, or a recorded G-I waiver from Vadim. Do not quietly proceed.",
           top + 1.86, height=1.00, tsize=12.5)
    two_col(
        sl,
        ("URL structure, two hierarchies of different depth", [
            "The homepage is the FitXpress parent, since /fitxpress/ redirects to /",
            "FitXpress vertical pages sit at /for-{vertical}/ with two-level breadcrumbs",
            "Mobile Tailor has its own parent, with /mobile-tailor/for-{vertical}/ beneath it and three-level breadcrumbs",
            "Never invent a /fitxpress/ path level",
        ], NAVY),
        ("Where things go", [
            "Command: /page [vertical] [gate | build | judge | handoff | full]",
            "Artifacts: workspace/pages/{slug}/",
            "quality-controller is not a substitute for G-J. It is neither blind nor page-shaped",
            "The benchmark to match: /structured-body-data-for-telehealth-digital-health-programs/",
        ], BLUE),
        top + 3.02, height=1.96, lfill=G50, rfill=BLUE_50, bsize=8.4, tsize=10.5,
    )


page(s_gates)


def s_slots(sl):
    top = slide_head(sl, "06 · LANDING-PAGE PLAYBOOK", "The section map, slots 1 to 9",
                     "The same 17 slots the pipeline enforces. Use this map when writing a brief, reviewing a draft, or building a campaign landing by hand.")
    table(sl, ["#", "Slot", "The rule"], [
        ["1", "Breadcrumbs", "Two levels for FitXpress (Home to Vertical), three for Mobile Tailor. Match what live pages declare. Never a middle level that resolves to a redirect"],
        ["2", "H1", "[Outcome in the vertical's terms] for [vertical]. One H1. Primary query in it and in the first 100 words. No \"best\", no \"most accurate\", no banned words"],
        ["3", "Hero and a vertical proof point", "One sentence on what the product does here, and a number from THIS vertical. The company-wide 112,100 scans figure does not work in a vertical hero"],
        ["4", "Vertical context", "The slot the page exists for. Three to five specifics only someone who has worked in this market knows: which regulators actually come up, who signs off and who blocks, the procurement cycle, the units, the export formats, the seasonality"],
        ["5", "Pains of this vertical", "From the use-case file and the audience.md hook, in the buyer's phrasing. Then honour that segment's \"what NOT to say\""],
        ["6", "What it is here, and the boundary", "Scope in the vertical's own words: what is captured, what is returned, what is documented. Then one boundary sentence: FitXpress is not positioned as a medical device. One negative, once"],
        ["7", "Where the workflow differs", "Not the whole flow. The two or three steps that run differently here: consent capture, retake logic, who reviews a flagged scan, how the record is filed"],
        ["8", "Compliance and data governance", "For a regulated vertical this is the deciding block. HIPAA and GDPR posture, encryption at rest, retention, no personal identifiers, consent handling, audit trails, privacy@3dlook.me. Every control stated with its limit"],
        ["9", "Accuracy, scoped", "The reframe, then the four conditions, then figures from proof-points.md only. No bare percentage. Never the three reserved words"],
    ], top, widths=[0.035, 0.225, 0.74], row_h=0.50)


page(s_slots)


def s_slots2(sl):
    top = slide_head(sl, "06 · LANDING-PAGE PLAYBOOK", "The section map, slots 10 to 17",
                     "Campaign landing variant: the same spine, shorter. Slots 2, 3, 5, 6, 9, a FAQ trimmed to four or five questions, and 15. One message, one action, no sibling-vertical cards competing with the campaign.")
    table(sl, ["#", "Slot", "The rule"], [
        ["10", "Cases from this vertical only", "Two cards minimum, each with a number from case-studies/, linked to /case-studies/. A case from an adjacent vertical breaks the page's promise. Mobile Tailor customer annual recurring revenue never appears anywhere"],
        ["11", "Customer quote from this vertical", "Name, role, company, only where the use is approved. No approved quote means the slot is dropped and recorded. Never an invented name or testimonial"],
        ["12", "Integration, formats, support", "The interface, web and mobile kits, the web widget, CSV export, 3D model export, what it plugs into here, white-label options, implementation support. The signal is \"we are already inside your environment\""],
        ["13", "Vertical FAQ", "The questions asked here and nowhere else. Ships with FAQPage schema"],
        ["14", "Price signal", "Name the entry tier and link to /pricing/. Mention the trial where it applies. Never the internal per-request rates"],
        ["15", "Primary action and form", "One action in the site's own language, visible without scrolling and repeated at the end. Minimal fields, visible consent, a confirmation state"],
        ["16", "Soft alternative and sibling verticals", "For buyers not ready to talk: the accuracy framework article, the vertical's hub article, a checklist. Then cards for two sibling verticals and a link up to the parent"],
        ["17", "Hidden technical layer", "Service or Product schema with audience and areaServed, FAQPage on the FAQ block, BreadcrumbList on the crumbs, canonical to self, Yoast title 60 characters or fewer, description 155 or fewer, a clean URL"],
    ], top, widths=[0.04, 0.235, 0.725], row_h=0.475)
    txt(sl, ML, top + 4.10, CW, 0.40,
        "Design compliance on a page: navy carrying 60 to 70% of the weight on hero, proof band, CTA band "
        "and footer, white on content zones, electric blue as a single sharp accent, Satoshi only, the "
        "eyebrow label above each section heading, cards at 20px and chips at 15px, oversized numerals in "
        "the proof zone. The comparison table reflows on mobile, or it is not shippable.",
        size=8.4, color=G600, italic=True, line_spacing=1.3)


page(s_slots2)


def s_faq(sl):
    top = slide_head(sl, "06 · LANDING-PAGE PLAYBOOK", "The FAQ block",
                     "The FAQ is where a landing page earns citations in answer engines, and it is the slot most often skipped. Only one page on the whole site currently ships FAQPage schema.")
    cards(sl, [
        ("Ten to thirteen questions",
         "The benchmark telehealth page runs 13. Below about eight, the block stops carrying weight."),
        ("Answer-first",
         "Each answer opens with a 40 to 60 word capsule that answers the question completely on its "
         "own, then adds detail. An answer engine quotes the first sentences."),
        ("Number first where a number exists",
         "\"Structured body data returns in about 45 seconds from two guided smartphone photos\" beats "
         "\"Results are fast.\""),
        ("Vertical questions only",
         "Regulatory posture, licensing, data residency, retention, who owns the data, what happens on a "
         "failed scan, what a pilot looks like, what the integration costs the development team. General "
         "product questions belong on the product page."),
        ("Source them from real objections",
         "From faq.md and from the business development owner for that market, not from imagination. "
         "Nick's call transcripts in Drive are the best raw material."),
        ("The boundary questions are mandatory",
         "On every health vertical: \"Is it a medical device?\" and \"Where do the photos go?\". The first "
         "answers no, structured body data supporting workflow, not a diagnostic device or a replacement "
         "for clinical judgment. The second gives the flow, the deletion policy and who controls retention."),
    ], top, cols=3, height=1.44, tsize=10.5, bsize=8.4, accent_bar=True)
    banner(sl, "FAQPAGE SCHEMA ON THE BLOCK",
           "Modelled on the telehealth page. Mark it up, or the block is doing half its job. Two more "
           "page-level rules help humans and answer engines equally: open every major section with a 40 to "
           "60 word answer-first capsule, and keep a visible \"Last updated\" line.",
           top + 3.14, height=0.94, tsize=12)


def s_cta(sl):
    top = slide_head(sl, "06 · LANDING-PAGE PLAYBOOK", "Calls to action, by funnel stage",
                     "Match the call to action to where the page sits. Do not force one call to action onto every page.")
    table(sl, ["Stage", "Page intent", "Call-to-action language"], [
        ["TOFU", "Educational, definitional, \"what is this\"",
         "Soft: \"Explore how mobile body scanning works\", \"Read the accuracy framework\""],
        ["MOFU", "Comparison, workflow, evaluation",
         "Evaluation: \"See how FitXpress supports remote progress tracking\", \"Review the buyer checklist\", \"See sample outputs\""],
        ["BOFU", "Operational, commercial, vertical page bottom",
         "Direct: \"Book a demo\", \"Talk to 3DLOOK about your workflow\", \"Explore FitXpress for telehealth and weight loss\""],
    ], top, widths=[0.09, 0.28, 0.63], row_h=0.42, bold_first=True)
    two_col(
        sl,
        ("Pair a primary with an evaluation secondary", [
            "On a vertical page, pair a direct primary with an evaluation secondary, which is what the benchmark page does: \"Book a demo\" plus \"See sample outputs\"",
            "The secondary catches the buyer who is not ready to talk to a human but is ready to look at a payload, and that buyer is a large share of the technical audience",
            "One primary action, visible without scrolling and repeated at the end",
            "Form fields minimal, consent visible, a confirmation state",
        ], BLUE),
        ("Site-wide language already in use, worth matching", [
            "\"Let's talk\"",
            "\"Request a Demo\"",
            "\"Unlock Body Data\"",
            "\"Get in touch\"",
            "\"Book a Consultation\"",
            "\"Explore Technology\"",
        ], NAVY),
        top + 1.90, height=1.90, lfill=BLUE_50, rfill=G50, bsize=8.8, tsize=11,
    )
    txt(sl, ML, top + 3.98, CW, 0.36,
        "Design compliance on a page: navy carrying 60 to 70% of the weight on hero, proof band, CTA band "
        "and footer, white on content zones, electric blue as a single sharp accent, Satoshi only, the "
        "eyebrow label above each section heading, cards at 20px and chips at 15px, oversized numerals in "
        "the proof zone. The comparison table reflows on mobile, or it is not shippable.",
        size=8.6, color=G600, italic=True, line_spacing=1.3)


page(s_faq)
page(s_cta)



def s_prelaunch(sl):
    top = slide_head(sl, "06 · LANDING-PAGE PLAYBOOK", "Page pre-launch checklist",
                     "Run it in order. Anything unresolved goes into the open items block rather than into a judgement call made alone.")
    cards(sl, [
        ("Structure and numbers", [
            "One H1, with H2 and H3 hierarchy intact",
            "Every number byte-identical across hero, body, FAQ and disclaimer",
            "Data-flow wording accurate everywhere",
            "Word count in the benchmark range, roughly 1,200 to 1,600",
            "A visible \"Last updated\" line",
        ]),
        ("Permissions and price", [
            "Logo-usage permission confirmed per customer",
            "No scan-volume figures published from the activity dashboard",
            "No pricing table on the page. Only the signal and a link to /pricing/",
            "The entry tier named, and the trial mentioned where it applies",
        ]),
        ("Links and markup", [
            "Internal links in all four directions: up to the parent, down or across to siblings, out to the supporting article, in from the parent",
            "JSON-LD present: Organization, Service or SoftwareApplication, BreadcrumbList, FAQPage",
            "Canonical to self, Yoast title 60 characters or fewer, description 155 or fewer",
        ]),
        ("Media and performance", [
            "Alt text written for this vertical's context",
            "WebP or AVIF, lazy-loaded",
            "Largest contentful paint under 2.5 seconds",
            "The comparison table reflows on mobile",
        ]),
        ("Sign-off", [
            "Asselya's editorial review before publish",
            "G-T passed",
            "G-J scored 85 or above, or the shortfall explicitly flagged",
            "Open items block written",
        ]),
    ], top, cols=5, height=2.72, tsize=10.5, bsize=8.2, accent_bar=True)
    banner(sl, None,
           "Asselya's editorial review and the G-J score are the two steps most likely to be skipped "
           "under deadline, and they are the two that keep a page from being republished twice.",
           top + 2.92, height=0.70, tsize=12)


page(s_prelaunch)


# ============================================================ 07 MESSAGING
def s_div7(sl):
    divider(sl, "07", "Messaging and proof library",
            "The approved hero lines, the whole permitted universe of figures, the citation rules, and the objections.")


page(s_div7, dark=True)


def s_heroes(sl):
    top = slide_head(sl, "07 · MESSAGING AND PROOF", "Approved hero lines",
                     "From messaging.md, which is the approved set. Adapt wording per artifact. Never invent a new claim inside a hero line.")
    table(sl, ["Segment", "Hero line"], [
        ["Telehealth and weight loss", "Verify body progress remotely to boost retention, reduce drop-off, and prove program ROI"],
        ["Online pharmacy, BMI verification", "Verify BMI inside the order flow to cut fraud, speed approvals, and stay audit-ready"],
        ["Insurance underwriting", "Verify body metrics remotely to issue faster, cut rework, and strengthen auditability"],
        ["Wellness rewards", "Verify wellness progress remotely to reduce disputes, boost participation, and improve program reporting"],
        ["Bariatric pre-authorization", "Pre-qualify patients remotely to reduce wasted consults, speed pre-auth, and improve conversion to procedures"],
        ["Occupational health", "Standardize screening intake remotely to increase throughput, reduce rescreens, and speed clearance decisions"],
        ["Clinical trials", "Standardize anthropometrics across sites and reduce visit burden to improve data quality and retention"],
    ], top, widths=[0.28, 0.72], row_h=0.375, bold_first=True)
    txt(sl, ML, top + 2.94, CW, 0.44,
        "Three of these are three-part parallel constructions, which sits close to the punch-triad rule. "
        "They are approved as hero lines because the three items are outcomes rather than adjectives. Do "
        "not extend the pattern into body copy.",
        size=9, color=G600, italic=True, line_spacing=1.28)
    two_col(
        sl,
        ("Taglines, used sparingly when a short hook is needed", [
            "Verified body data, built for trust",
            "Two photos. 80+ measurements. 45 seconds",
            "From scan to outcome, in one workflow",
        ], BLUE),
        ("Phrases to prefer", [
            "\"Verified body data\" over \"accurate measurements\"",
            "\"Trusted workflow layer\" over \"scanning API\"",
            "\"Audit-ready records\" for regulated industries",
            "\"Real-world accuracy\" over \"lab accuracy\"",
        ], NAVY),
        top + 3.44, height=1.50, lfill=BLUE_50, rfill=G50, bsize=8.8, tsize=11,
    )


page(s_heroes)


def s_proof(sl):
    top = slide_head(sl, "07 · MESSAGING AND PROOF", "The approved proof set",
                     "This is the whole permitted universe of figures, copied from proof-points.md. If a number is not here, it does not ship.")
    stat_tiles(sl, [
        ("96 to 97%", "Accuracy against expert manual measurement"),
        ("1.5 to 2.0 cm", "Typical error margin"),
        ("< 1 cm", "Variance across repeated scans"),
        ("95%+", "Overall repeatability consistency"),
        ("+/- 3.5%", "Weight estimation average error"),
        ("< 45 sec", "Two photos to structured results"),
    ], top, height=1.10, num_size=15.5, lab_size=8.0)
    y = top + 1.24
    txt(sl, ML, y, CW * 0.485, 0.22, "ACCURACY BY MEASUREMENT, ABSOLUTE ERROR",
        size=8, color=BLUE, bold=True)
    tw = CW * 0.485
    acc = [["Wrist girth", "0.54 cm"], ["Calf", "1.27 cm"], ["Neck", "1.48 cm"],
           ["Thigh", "1.64 cm"], ["Knee", "1.73 cm"], ["Chest", "1.74 cm"],
           ["Waist", "2.14 cm"], ["Hip", "2.25 cm"]]
    for j, r in enumerate(acc):
        ry = y + 0.28 + j * 0.245
        rect(sl, ML, ry, tw, 0.245, fill=G50 if j % 2 == 0 else WHITE)
        txt(sl, ML + 0.12, ry + 0.04, tw * 0.6, 0.2, r[0], size=8.4, color=G700,
            fit_check=False)
        txt(sl, ML + tw * 0.6, ry + 0.04, tw * 0.4 - 0.14, 0.2, r[1], size=8.4,
            color=BLUE, bold=True, align=PP_ALIGN.RIGHT, fit_check=False)
    x2 = ML + tw + 0.30
    tw2 = CW - tw - 0.30
    txt(sl, x2, y, tw2, 0.22, "REPEATABILITY BY MEASUREMENT, SAME STUDY", size=8,
        color=BLUE, bold=True)
    rep = [["Chest", "0.60 cm"], ["Waist", "0.89 cm"], ["Low hips", "0.86 cm"],
           ["Knee", "0.12 cm"], ["Calf", "0.12 cm"], ["Ankle", "0.07 cm"]]
    for j, r in enumerate(rep):
        ry = y + 0.28 + j * 0.245
        rect(sl, x2, ry, tw2, 0.245, fill=G50 if j % 2 == 0 else WHITE)
        txt(sl, x2 + 0.12, ry + 0.04, tw2 * 0.6, 0.2, r[0], size=8.4, color=G700,
            fit_check=False)
        txt(sl, x2 + tw2 * 0.6, ry + 0.04, tw2 * 0.4 - 0.14, 0.2, r[1], size=8.4,
            color=BLUE, bold=True, align=PP_ALIGN.RIGHT, fit_check=False)
    txt(sl, x2, y + 1.62, tw2, 0.80,
        "Both tables come from the 2025 Accuracy and Repeatability Study, measured against expert manual "
        "measurement. Never mix them with the ISO multi-company benchmark, which is a different reference "
        "and is not in proof-points.md. Repeatability is the argument for every longitudinal vertical, so "
        "lead with it there.",
        size=8.6, color=G600, italic=True, line_spacing=1.3)


page(s_proof)


def s_proof2(sl):
    top = slide_head(sl, "07 · MESSAGING AND PROOF", "The approved proof set, continued",
                     "Coverage, training data, customers and company facts. Same rule: if a number is not in proof-points.md, it does not ship.")
    cards(sl, [
        ("Speed and coverage",
         "Photo to results under 45 seconds. Photos required: 2, front and side. Body measurements: 80+. "
         "Body composition outputs: Body Mass Index (BMI), basal metabolic rate (BMR), body fat "
         "percentage, lean mass, fat mass, essential fat, beneficial fat. Points in the source 3D model: "
         "5M+ per model."),
        ("Training data",
         "9+ years. 150,000+ photographs. 30,000+ 3D scans. 430,000+ individual measurements. Ages 16 to "
         "78. Weight 38 to 210 kg. Height 150 to 220 cm. 48% male, 52% female. Locations: US and Europe. "
         "86 parameters measured per person, 34 photo configurations per user."),
        ("Company",
         "Founded 2016. 28 employees. $16.2M raised. Sifted 2020 Pioneers of the New World. IEEE Retail "
         "Digital Transformation Grand Challenge winner. Member of Mobile Body Scanning Standards. "
         "Headcount is contested: the sales playbook says 22. See section 9, item 8."),
        ("FitXpress customer outcomes",
         "Yazen: 34,000 scans in 2025, weight-loss management support. UK Meds: 7,500 scans in 2025, Body "
         "Mass Index (BMI) verification for online pharmacy, 7 months customer lifetime to date. Healthyr: "
         "patient profile complement. Naming is subject to permission, and permission is currently "
         "contradictory. See section 9, item 5."),
        ("Aggregate",
         "100+ customers all-time. 67 active customers in 2025. 112,100 total scans in 2025. Internal only "
         "and never customer-facing: annual recurring revenue figures, and Mobile Tailor customer revenue "
         "in any form."),
        ("Market sizing, internal only",
         "Insurance underwriting $25 to 75M a year. Wellness rewards $50 to 200M. Bariatric "
         "pre-qualification $10 to 30M. Occupational health screening $20 to 60M. Clinical trials $10 to "
         "40M. Marked illustrative in the source. Planning and prioritization use only."),
    ], top, cols=3, height=1.86, tsize=10.5, bsize=8.4, fill=BLUE_50, accent_bar=True)
    banner(sl, None,
           "A vertical hero needs a number from THAT vertical. The company-wide 112,100 scans figure does "
           "not work in a vertical hero, and reaching for it is the usual sign that the vertical has no "
           "proof yet.",
           top + 4.00, height=0.72, tsize=12)


page(s_proof2)


def s_compliance(sl):
    top = slide_head(sl, "07 · MESSAGING AND PROOF", "Compliance, price signal, and the citation rules",
                     "The compliance block is the deciding block in every regulated vertical, and the negatives in it are as load-bearing as the positives.")
    two_col(
        sl,
        ("Compliance and security, what is true", [
            "HIPAA maintained, for FitXpress in US healthcare contexts. Follows GDPR principles",
            "TLS in transit. AWS S3 server-side encryption at rest, always on, cannot be disabled",
            "Photos removed immediately after processing or within 30 days per client policy, and auto-blurred when stored. Face obfuscation at capture",
            "No personal identifiers processed. End-user images never shared with third parties",
            "Business associate agreements signed for HIPAA-covered customers. Privacy contact privacy@3dlook.me",
        ], BLUE),
        ("What is NOT true, and must never be implied", [
            "NOT cleared by the Food and Drug Administration (FDA)",
            "NOT SOC 2 certified. It is in progress. Confirm with Vadim before any mention",
            "NOT peer-reviewed, and NOT third-party clinically validated",
            "Compliance is framed on data-privacy frameworks, never on medical-device frameworks",
        ], NAVY),
        top, height=2.10, lfill=BLUE_50, rfill=G100, bsize=8.4, tsize=11,
    )
    cards(sl, [
        ("Pricing signal",
         "Free trial: one month, 200 requests, full kit access. Entry tier $1,000 a month. No integration "
         "fee, no setup fee. Public detail lives on /pricing/, and the internal per-request table never "
         "appears in an external artifact. Confirm the live page before any deck or landing page ships, "
         "because pricing.md and the live page disagree. Section 9, item 4."),
        ("Citation rules, part one", [
            "Prefer the customer outcome over the internal metric. \"UK Meds uses FitXpress for BMI verification\" carries more weight with a buyer than \"we have 96% accuracy\"",
            "Never invent a comparison. No \"ten times more accurate than X\" unless that exact figure exists in proof-points.md with a source",
            "Use the range form when a figure may drift. \"Around 95 to 97% accuracy\" survives a small correction; a single precise number does not",
        ]),
        ("Citation rules, part two", [
            "One number, everywhere the same. Byte-identical in the hero, the body, the FAQ and the disclaimer. If two sources conflict, never average them: keep the defensible one, or replace both with a qualitative statement",
            "Every customer name and number comes from case-studies/, and only with current written permission",
            "External statistics need a named body and a link. Verify at source, and get approval before it appears in a customer-facing deck",
        ]),
    ], top + 2.26, cols=3, height=2.14, tsize=10.5, bsize=8.2, accent_bar=True)


page(s_compliance)


def s_objections(sl):
    top = slide_head(sl, "07 · MESSAGING AND PROOF", "Objection handling",
                     "Marketing-side answers built for collateral. Sales has its own phrasing in the AE playbook, and the two should not contradict each other.")
    table(sl, ["Objection", "Answer"], [
        ["\"Patients already self-report their weight.\"",
         "Self report is inconsistent and easy to misrepresent. FitXpress adds structured, repeatable body data before review, from two photos in under 45 seconds, inside the flow the member already completes"],
        ["\"We already use connected scales.\"",
         "A scale returns one number. FitXpress returns body composition, 80+ measurements, a 3D model and scan-to-scan comparison, captured remotely from a phone"],
        ["\"We don't want onboarding friction.\"",
         "Two photos with guided capture, under a minute. The scan is triggered from a flow the program already runs, and the capture layer validates pose and framing in real time to reduce retakes"],
        ["\"How accurate is it?\"",
         "Accurate enough for which decision? Against expert manual measurement, 96 to 97% with a typical absolute error of 1.5 to 2.0 cm. For repeat-scan programs the figure that matters is repeatability: variance stays < 1 cm, and a change between scans reads as real rather than as noise. Detailed methodology is available under NDA"],
        ["\"Our clinicians and underwriters assess manually.\"",
         "The human still decides. FitXpress standardizes intake and reduces review burden, and it supports clinician or underwriter review rather than replacing it"],
        ["\"How does this help the business?\"",
         "Better qualification and engagement feed conversion, retention and cleaner outcomes reporting to the payers, employers and regulators the program answers to"],
    ], top, widths=[0.26, 0.74], row_h=0.52, bold_first=True)
    txt(sl, ML, top + 3.34, CW, 0.24,
        "Every one of these answers hedges, and none of them oversells. That is the point: an answer a "
        "buyer can check is worth more than an answer that wins the moment.",
        size=9, color=G600, italic=True)


page(s_objections)


def s_objections2(sl):
    top = slide_head(sl, "07 · MESSAGING AND PROOF", "The boundary objections, and competitive framing",
                     "These four come up in diligence rather than in discovery, and a weak answer here costs the deal later.")
    table(sl, ["Objection", "Answer"], [
        ["\"We're evaluating other options.\"",
         "Compare by role. The useful question is which method fits which workflow: a remote structured-data layer, an in-clinic reference method, and a consumer app answer different questions. Route named-competitor framing to Katerina before it goes in writing"],
        ["\"Is this a medical device?\"",
         "No. FitXpress is not positioned as a medical device. It provides structured body data that supports clinical and operational workflows, and compliance is evaluated on data-privacy frameworks rather than medical-device frameworks"],
        ["\"Do you train your AI on our photos?\"",
         "No. Photos sent through a customer's tenant are deleted per the retention policy and are not used to train the model"],
        ["\"Where is the data stored, and for how long?\"",
         "AWS S3 with mandatory server-side encryption, TLS in transit. Photos are removed immediately after processing or within 30 days, depending on the client policy chosen, and are auto-blurred when stored. No personal identifiers are processed"],
    ], top, widths=[0.26, 0.74], row_h=0.54, bold_first=True)
    txt(sl, ML, top + 2.50, CW, 0.24, "COMPETITIVE FRAMING", size=8.6, color=BLUE,
        bold=True)
    cards(sl, [
        ("The landscape",
         "Prism Labs is the primary competitor in the FitXpress space, strong in insurance, population "
         "health and GLP-1. Bodygram is secondary, aimed at trainers and dieticians. Size Stream is strong "
         "in hybrid on-premise and at-home clinical research. The long-term risk is native platform "
         "primitives from Apple or Google."),
        ("Three rules",
         "Compare by method or by role, never by competitor name in published copy. Never a clean sweep, "
         "because a table where we win every row reads as marketing and both buyers and answer engines "
         "discount it. Route any competitive positioning to Katerina before it goes in writing, because "
         "competitive data moves and the internal analysis is dated March and April 2026."),
        ("The internal boundary",
         "The competitive analysis deck in Drive is internal only and never leaves the building. Nothing "
         "from it is quotable in a customer-facing artifact, including in a comparison table that does not "
         "name the competitor."),
    ], top + 2.78, cols=3, height=1.90, tsize=10.2, bsize=8.2, accent_bar=True)


page(s_objections2)


# ============================================================ 08 QA
def s_div8(sl):
    divider(sl, "08", "QA checklist",
            "Run this before anything is sent, published or handed to a designer. It maps to the 20-point rubric in docs/quality-rubric.md.")


page(s_div8, dark=True)


def s_qa1(sl):
    top = slide_head(sl, "08 · QA CHECKLIST", "Claims and facts, and voice and brand",
                     "Rubric categories B and C. The claims list is the hard-fail category: a single miss here is worse than a whole deck of weak copy.")
    two_col(
        sl,
        ("Claims and facts", [
            "Every figure traced to a line in proof-points.md. No exceptions, no \"approximately\"",
            "Every figure byte-identical across hero, body, FAQ, disclaimer and any chart label",
            "Every customer name and number from case-studies/, with current written permission",
            "No Mobile Tailor content or customer revenue anywhere",
            "No claim of diagnosis, treatment, eligibility, underwriting, hiring or clearance decisioning",
            "No claim of replacing a clinician, DEXA, BIA, a calibrated scale or a protocol reference method",
            "No guaranteed compliance, no automatic fraud detection, no \"most accurate\", no \"just an API\"",
            "Accuracy scoped: the reframe, the four conditions, no bare headline percentage",
            "Repeatability written as < 1 cm and never presented as accuracy. The two benchmarks are not mixed",
            "\"Independent\", \"validated\", \"third-party\" absent unless a named external party and a citable output are present",
            "SOC 2 not claimed. FDA clearance not claimed. Peer review not claimed",
            "Roadmap items limited to shipped or in-progress, with no promised dates",
            "External statistics carry a named body and a link, and have external-use approval",
        ], BLUE),
        ("Voice and brand", [
            "Banned-word grep clean, across all 20 words including figurative \"navigate\"",
            "No em dash or en dash anywhere",
            "No \"it's not just X, it's Y\", no \"not only X but also Y\"",
            "No adjectival punch triad",
            "Terminology guardrails clean: no \"objective\" about our output, no \"the reader\", no \"this guide\", no \"by hand\", no \"plus\" as a connector, no \"let\", no \"so\" introducing a benefit, no corrective \"X, not Y\" outside a stated boundary",
            "M1: every abbreviation expanded at first use, including BMI, GLP-1, API, SDK and every cited regulator",
            "M2: one clear negative scope statement, with no second negation chained onto it",
            "The hero sells an outcome, not accuracy",
            "Buyer framing rather than \"you\"-spam, except on conversion sections",
            "The detector has been run and every hard fail and house-rule violation is fixed",
            "The second pass has been done: what here still reads as machine-written, and it was fixed",
        ], NAVY),
        top, height=4.30, lfill=BLUE_50, rfill=G50, bsize=8.4, tsize=12,
    )


page(s_qa1)


def s_qa2(sl):
    top = slide_head(sl, "08 · QA CHECKLIST", "Disclaimers, design, and conversion",
                     "Three columns for three different readers. Hand a designer the design column, not the whole list.")
    cards(sl, [
        ("Disclaimers and sensitive verticals", [
            "The vertical's scope note appears early, rather than only at the end",
            "The italic disclaimer sits near any accuracy or eligibility claim",
            "\"FitXpress is not positioned as a medical device\" is present, in that wording",
            "The segment's \"what NOT to say\" list from section 4 checked line by line",
            "Every control stated with its limit",
            "BCRL, occupational health, clinical trials and insurance drafts have gone to Whitney",
        ]),
        ("Design compliance", [
            "#143DFF as a single sharp accent, never a large fill. #2962FF absent",
            "Navy #050F40 with the radial glow and texture on hero, proof, CTA and footer. No flat navy",
            "Satoshi only. No Inter, Bricolage Grotesque or IBM Plex Sans",
            "Type scale, spacing steps and radii from DESIGN.md, with nothing off-scale",
            "Eyebrow label on each section or content slide",
            "Proof zone leads with an oversized numeral",
            "Product imagery in preference to icon grids",
            "AA contrast, a scrim under light text on imagery, a visible focus ring, reduced-motion honoured",
            "White call-to-action button on navy",
        ]),
        ("Calls to action and conversion", [
            "The CTA matches the funnel stage, and one CTA has not been forced onto everything",
            "One primary action, visible without scrolling and repeated at the end",
            "An evaluation-stage secondary action exists for the buyer who is not ready to talk",
            "The price signal names the entry tier and links to /pricing/. No internal per-request rates",
            "Form fields minimal, consent visible, a confirmation state exists",
        ]),
    ], top, cols=3, height=4.30, tsize=11, bsize=8.5, accent_bar=True)


page(s_qa2)


def s_qa3(sl):
    top = slide_head(sl, "08 · QA CHECKLIST", "Structure, and sign-off",
                     "The structure column splits: pages first, then the deck-only items. Sign-off is five names and two gates, and none of them are optional.")
    cards(sl, [
        ("Structure: pages", [
            "One H1, hierarchy intact. Ten or more FAQ questions, answer-first, with FAQPage schema",
            "Service or Product schema with audience and areaServed, BreadcrumbList, canonical to self",
            "Yoast title 60 characters or fewer, description 155 or fewer, clean URL",
            "Breadcrumb depth matches the real hierarchy, with no invented path level",
            "Internal links in all four directions. Comparison table reflows on mobile",
            "A visible \"Last updated\". Word count roughly 1,200 to 1,600",
        ]),
        ("Structure: decks", [
            "The 19-slide core is present, and any dropped slide is a recorded decision",
            "Fixed slides copied verbatim from the last approved deck, not rewritten",
            "The text-fit review is written, slide by slide",
            "The confidentiality footer carries the current month and year",
            "Empty slots, the case card and the testimonial, are recorded rather than filled with adjacent-vertical material",
        ]),
        ("Sign-off", [
            "Open items block written. Per guardrail 11, unresolved trade-offs are flagged rather than silently decided",
            "Numbers confirmed by Vadim",
            "Regulated claims reviewed by Whitney",
            "Competitive framing reviewed by Katerina",
            "Editorial review by Asselya complete",
            "For a site page: G-T passed and G-J scored 85 or above, or the shortfall explicitly flagged",
        ]),
    ], top, cols=3, height=4.30, tsize=11, bsize=8.5, accent_bar=True)


page(s_qa3)


# ============================================================ 09 OPEN ITEMS
def s_div9(sl):
    divider(sl, "09", "Open items and conflicts",
            "Found while building this playbook. Per guardrail 11, these are flagged rather than silently decided. The first audit starts here and extends the list.")


page(s_div9, dark=True)


def s_open1(sl):
    top = slide_head(sl, "09 · OPEN ITEMS", "The blockers, ranked",
                     "The top four decide what marketing can honestly promise. Everything below them is a correction rather than a decision.")
    items = [
        ("02", "The ISO benchmark set is in circulation and is not in proof-points.md",
         "0.40 cm session-to-session repeatability, ISO 8559-1:2017, 14 companies, 8 countries, 1,152 data "
         "points, and the comparison figures 0.57 cm for 3D scanners and 0.94 cm for expert manual "
         "measurement. All of it appears in about-me.md, the sales playbook and the shipped insurance "
         "deck. None of it is in proof-points.md, so under the sourcing rule it is unusable, which "
         "conflicts with two shipped assets. ASK: add it with its source, or withdraw it from collateral. "
         "The highest-value item here, because 0.40 cm is the strongest number we have for longitudinal "
         "verticals."),
        ("12", "G-I blocks nearly every FitXpress vertical page",
         "The gate requires two or more publishable cases from the vertical, and every FitXpress vertical "
         "has at most one. ASK Vadim for a second case per priority vertical, an approved external "
         "reference, or a recorded waiver. Without this, \"vertical landing pages for all use cases\" is "
         "not deliverable as stated."),
        ("05", "Customer naming permission is contradictory",
         "The insurance deck's open items record an NDA prohibiting naming Yazen, UK Meds and Healthyr. "
         "The telehealth page structure document plans an eleven-logo wall. The sales playbook names three "
         "clients on a competitive slide. case-studies/ tells agents to cite UK Meds and Yazen by name in "
         "outbound. ASK Katerina and Olena Chorna for a single per-customer permission list. Until it "
         "exists, treat naming as blocked and use the numbers without the name."),
        ("04", "Pricing figures disagree across three sources",
         "pricing.md says $1,000 a month for 500 requests and $1,500 for 1,000. The sales playbook says "
         "$1,000 for 1,000 requests with a $0.50 per-request floor above 5,000 a month, where pricing.md "
         "puts $0.50 at 20,000. There is also a live /pricing/ page found to contradict pricing.md. "
         "Separately, about-me.md says never state or imply prices, while the page pipeline requires a "
         "price signal on every commercial page. ASK: which table is current, and what the standing rule "
         "is."),
    ]
    cols, gap = 2, 0.24
    w = (CW - gap) / 2
    h = 1.96
    for i, (n, t, b) in enumerate(items):
        r, c = divmod(i, cols)
        x = ML + c * (w + gap)
        y = top + r * (h + 0.20)
        rect(sl, x, y, w, h, fill=G50, radius=20)
        rect(sl, x, y, 0.055, h, fill=BLUE, radius=None)
        txt(sl, x + 0.26, y + 0.16, 0.6, 0.3, n, size=17, color=BLUE_200, bold=True,
            fit_check=False)
        txt(sl, x + 0.86, y + 0.20, w - 1.10, 0.44, t, size=10.8, color=NAVY,
            bold=True, line_spacing=1.14)
        txt(sl, x + 0.26, y + 0.72, w - 0.50, h - 0.90, b, size=8.2, color=G700,
            line_spacing=1.26)


page(s_open1)


def s_open2(sl):
    top = slide_head(sl, "09 · OPEN ITEMS", "The corrections, and the gaps",
                     "Items 1, 3, 6, 7, 8, 9, 10, 11, 13, 14, 15 and 16. Most are one confirmation away from resolved.")
    table(sl, ["#", "Item", "The ask"], [
        ["01", "#2962FF named as the accent in this playbook's own brief",
         "A superseded placeholder from before the Figma export. This playbook uses #143DFF and Satoshi. If #2962FF is genuinely in use somewhere, that artifact needs correcting rather than the design system"],
        ["03", "ISO 8559-1:2017 versus ISO 20685-1",
         "about-me.md, the sales playbook and the insurance deck cite one. The June 2026 telehealth page structure document cites the other, with a different benchmark shape. One is wrong and both are in circulation. Ask Vadim and product which standard the benchmark ran against"],
        ["06", "The speed claim varies: 45 seconds, under a minute, under 60 seconds, 40 to 50 seconds",
         "proof-points.md and tech-spec.md say under 45 seconds. The published 2-pager says under 60 seconds twice. The insurance deck cover says under a minute. Guardrail 2 makes this a defect. Hold \"under 45 seconds\" as the single claim, and correct the 2-pager"],
        ["07", "Training-data height range disagrees",
         "RESOLVED 2026-09-02 - Vadim confirmed 150 to 220 cm. proof-points.md, how-it-works.md and faq.md corrected from the old 150 to 205 cm deck figure"],
        ["08", "Company headcount disagrees",
         "CLAUDE.md and proof-points.md say 28 employees. The June 2026 sales playbook says 22. The About slide in every deck carries this number"],
        ["09", "Banned words are live in shipped assets",
         "\"Best-in-class Repeatability\" titles slide 6 of the insurance deck, and the phrase also sits in icp-detail.md and the sales playbook. The 2-pager uses \"seamlessly embeds\". Previous audits found \"leverage\", \"best-in-class\" and \"revolutionize\" on live site pages. Week-one audit produces the full list and a correction plan"],
        ["10", "\"Independently benchmarked\" on insurance slide 6",
         "Breaks guardrail 3, since independence is not provable with a named external party and a citable output. Compliant rewrite: \"benchmarked against 3D scanners and expert manual measurements in a multi-company benchmark\", with the naming resolved through item 2"],
        ["11", "Two verticals have no proof and no use-case file",
         "Plastic surgery and BCRL. Neither can carry a deck or a page yet. Ask: should product produce a volumetric-asymmetry validation figure, and is plastic surgery a 2026 priority?"],
        ["13", "FX Guidance for Decks, Whitney, February 2026, did not resolve",
         "A shortcut in the Drive Marketing docs folder whose target could not be read. Likely the single most relevant existing document for deck work. Ask Whitney or Vadim to re-share it"],
        ["14", "Digital fitness has no citable customer",
         "The names in circulation appear in the sales deck and the telehealth page draft but not in proof-points.md or case-studies/. Either add them with permission, or the fitness deck ships without a case card"],
        ["15", "HIPAA applicability to life insurers",
         "Flagged in the insurance deck's own open items, since HIPAA covers health plans and providers and a life insurer's status varies. The compliance slide carries the claim forward from the telehealth deck. Ask Whitney"],
        ["16", "The Munich Re and CDC statistics on insurance slides 2 and 7",
         "Carried from an article, and the deck flags them as needing external-use approval. Verify at source and get approval before the deck goes out again"],
    ], top, widths=[0.035, 0.265, 0.70], row_h=0.355)


page(s_open2)


# ---- close
def s_close(sl):
    navy_bg(sl, center=(50000, 26000))
    txt(sl, ML, 1.86, 10.4, 1.3, "You have the guardrails.\nNow ship one page.",
        size=38, color=WHITE, bold=True, line_spacing=1.1)
    rect(sl, ML, 3.42, 1.7, 0.05, fill=BLUE, radius=None)
    txt(sl, ML, 3.72, 8.4, 0.5,
        "Week two is the /for-bmi-verification/ rewrite. It is the weakest live page and the vertical "
        "with a real customer reference, so it exercises every gate in section 6.",
        size=12.5, color=BLUE_100, line_spacing=1.34)
    who = [
        ("Numbers and assets", "Vadim Bilan"),
        ("Regulated claims", "Whitney Cathcart"),
        ("Positioning and competitive", "Katerina Galich"),
        ("Editorial sign-off", "Asselya Sekerova"),
        ("Capability and roadmap", "Max Kucherenko"),
    ]
    gap = 0.18
    w = (CW - gap * 4) / 5
    for i, (lab, name) in enumerate(who):
        x = ML + i * (w + gap)
        rect(sl, x, 4.66, w, 0.92, fill=None, radius=20, line=BLUE_700, line_w=1.1)
        txt(sl, x + 0.18, 4.82, w - 0.36, 0.2, lab, size=7.6, color=BLUE_300,
            bold=True, fit_check=False)
        txt(sl, x + 0.18, 5.06, w - 0.36, 0.4, name, size=11.5, color=WHITE,
            bold=True, line_spacing=1.14)
    txt(sl, ML, 6.06, CW, 0.28,
        "Section 9 is the starting point for the week-one audit, not the end of it. "
        "Add to it rather than working around it.",
        size=10, color=BLUE_200, italic=True)
    txt(sl, ML, 6.62, CW, 0.24,
        "FitXpress Product Marketing Playbook  ·  Version 1.0  ·  August 2026  ·  "
        "Confidential, internal use only",
        size=8.5, color=BLUE_300, bold=True)


page(s_close, dark=True)


# ----------------------------------------------------------------------------- render
TOTAL = len(S)
for i, (fn, dark) in enumerate(S):
    sl = new_slide(prs)
    fn(sl)
    if i > 0:
        footer_line(sl, i + 1, TOTAL, dark=dark)

out = "/home/vadim_prod/3dlook-marketing/marketing_vb/workspace/playbooks/FitXpress-Product-Marketing-Playbook-Nika.pptx"
prs.save(out)

print(f"slides: {TOTAL}")
print(f"saved:  {out}")
if warnings:
    print(f"\n{len(warnings)} fit warnings:")
    for w in warnings:
        print(w)
else:
    print("\nno fit warnings")
